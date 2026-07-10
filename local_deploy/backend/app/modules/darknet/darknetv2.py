import sys
if "--osint-mapper" in sys.argv:
    run_osint_mapper()
    sys.exit(0)
"""
 - All-in-one OSINT & darknet crawler (Enhanced Edition v5.1 - OSINT Recon)

Features:
- Auto-create venv and install dependencies if needed, re-executes inside venv
- Auto-start Tor (Windows: tor.exe, Linux/macOS: tor) and verify via check.torproject.org
- Multi-DB support: MongoDB (URI), SQLite, MySQL, PostgreSQL (Patched Connection Handling)
- Recursive multithreaded crawling with query-parameter expansion
- 🧠 UIE REGEX ENGINE (v1.1): Robust entity extraction with Heuristic Context Filtering
- 🧠 NEMESIS INTELLIGENCE ONTOLOGY v3.0: Unified Entity Matrix & Autonomous Task Framework
- 📄 DOCUMENT PARSING ENGINE: Extracts text streams from PDF, DOCX, PPTX, CSV, TXT
- Real-time insert into DB + append newline-delimited JSON (export.jsonl)
- Real-time auto-export to CSV for extracted entities
- 🚀 ADVANCED RICH UI: Modern Interactive Console Tables & Telemetry
- Autonomous Subprocess OSINT Mapper (Auto-opens in secondary console with Auto-RECON)
"""

import os
import sys
import time
import json
import csv
import re
import hashlib
import random
import threading
import subprocess
import webbrowser
import socket
import concurrent.futures
import warnings
import io
import smtplib
import urllib3
import tempfile
from typing import List, Dict, Any, Optional
from email.mime.text import MIMEText
from collections import deque
from queue import Queue, Empty
from datetime import datetime, timezone
from urllib.parse import urlparse, urljoin, parse_qs, urlencode

# --- NEMESIS TRACER ENGINE ---
try:
    from nemesis_tracer import NemesisTracerEngine
    tracer_engine = NemesisTracerEngine()
except ImportError:
    tracer_engine = None
    # print("⚠️ Failed to import Nemesis Tracer Engine: No module named 'nemesis_tracer'")

# Suppress insecure request warnings if fetching without verify=False
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# Suppress BeautifulSoup XML parsing warnings for general web crawling
from bs4 import XMLParsedAsHTMLWarning
warnings.filterwarnings("ignore", category=XMLParsedAsHTMLWarning)

# -------------------- VENV BOOTSTRAP & DEPENDENCIES --------------------
VENV_DIR = os.path.join(os.getcwd(), "venv")
IN_VENV_FLAG = "INTEL_CRAWLER_IN_VENV"
REQUIRED_PKGS = [
    "requests[socks]>=2.31.0",
    "beautifulsoup4>=4.12.2",
    "pymongo>=4.4.0",
    "psycopg2-binary>=2.9.7",
    "mysql-connector-python>=8.0.32",
    "sqlalchemy>=2.0.20",
    "lxml>=4.9.3",
    "fake_useragent>=1.1.0",
    "python-dotenv>=1.0.0",
    "rich>=13.5.0",
    "stem>=1.8.1",
    "flask>=2.0.0",
    "flask-cors>=3.0.10",
    "pypdf>=3.17.0",
    "python-docx>=1.1.0",
    "python-pptx>=0.6.22",
    "google-api-python-client>=2.100.0",
    "google-auth-httplib2>=0.1.1",
    "google-auth-oauthlib>=1.1.0",
    "google-genai>=0.1.0",
    "playwright>=1.39.0",
    "openpyxl>=3.1.2"
]


# Dependencies managed centrally by main.py Orchestrator
# -------------------- IMPORTS (after venv ensured) --------------------
import requests
import pypdf
import docx
import pptx
import openpyxl
from playwright.sync_api import sync_playwright
from bs4 import BeautifulSoup
from fake_useragent import UserAgent
from dotenv import load_dotenv

# API & UI Imports
import flask.cli
from flask import Flask, request, jsonify, Response, render_template
from flask_cors import CORS
from rich.console import Console
from rich.live import Live
from rich.layout import Layout
from rich.panel import Panel
from rich.table import Table
from rich.text import Text
from rich.tree import Tree
from rich.align import Align
from rich import box

import os
gs_client = None
gs_sheet = None

def upload_to_google_sheets(record):
    pass

# DB drivers
try:
    from pymongo import MongoClient
except Exception:
    MongoClient = None
try:
    import sqlite3
except Exception:
    sqlite3 = None
try:
    import mysql.connector as mysql_connector
except Exception:
    mysql_connector = None
try:
    import psycopg2
except Exception:
    psycopg2 = None

load_dotenv()
console = Console()

# -------------------- CONFIG & GLOBALS --------------------
AUTONOMOUS_MODE = os.getenv("VITE_CRAWLER_ENABLED", "false").lower() == "true"
USE_TOR = os.getenv("VITE_TOR_AUTO_START", "true").lower() == "true"
TOR_PORT = os.getenv("VITE_TOR_SOCKS_PORT", "9050")
TOR_CHECK_URL = "http://check.torproject.org"
TOR_PROXY = {"http": f"socks5h://127.0.0.1:{TOR_PORT}", "https": f"socks5h://127.0.0.1:{TOR_PORT}"}

HEADERS = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
_thread_env = os.getenv("VITE_CRAWLER_MAX_THREADS", "10")
if _thread_env.startswith("${") or not _thread_env.isdigit():
    MAX_WORKERS = 10
else:
    MAX_WORKERS = int(_thread_env)
SAVE_STATE_INTERVAL = 30
STATE_FILE = "state.json"
EXPORT_JSONL = "export.jsonl"
EXPORT_CSV = "entities_export.csv"
SEED_FILE = "darknet_urls.txt"
KEYWORDS_FILE = "keywords.txt"
PLATFORMS_FILE = "platforms.json"

DEFAULT_KEYWORDS = ["bitcoin", "ethereum", "ransomware", "login verify", "account support", "private key"]

KNOWN_ENTITIES = {
    "kaspa:qp7r644htq3xramj68tzz7ndk95hz9mhuzq6rfh0tnnxzppss2sxxseu0ysfe": "MEXC_Hot_Wallet",
    "kaspa:qqe3p64wpjf5y27kxppxrgks2pwsnv3cvsm2pas2qcpgpd4gzzgw2twzcqmxs": "KuCoin_Hot_Wallet",
    "0x28c6c06298d514db089934071355e5743bf21d60": "Binance 14 (Cold)",
    "0xd90e2f925da726b50c4ed8d0fb90ad053324f31b": "Tornado Cash Router (MIXER)",
    "0xdf9b4b57865b403e08c85568442f95c26b7896b0": "Stargate Finance (BRIDGE)",
    "0x401f6c983ea34274ec46f84d70b31c15146b1f29": "Polygon POS Bridge (BRIDGE)",
}

# 🧠 NEMESIS INTELLIGENCE ONTOLOGY v3.0
NEMESIS_ONTOLOGY = {
    "PERSON": {
        "tasks": ["T1: Identity verification", "T2: Alias resolution", "T3: Social graph mapping", "T5: Phone ↔ email ↔ wallet linking", "T6: Risk scoring"]
    },
    "ORGANIZATION": {
        "tasks": ["T9: Corporate registry lookup", "T10: Director/shareholder mapping", "T11: Domain/IP linkage", "T13: Fraud cluster detection"]
    },
    "BANK_ACCOUNT": {
        "tasks": ["T15: Inflow/outflow tracing", "T16: Transaction clustering", "T18: AML risk scoring", "T20: Fiat → crypto bridge detection"]
    },
    "CRYPTO_WALLET": {
        "tasks": ["T21: Transaction graph building", "T22: Multi-hop tracing", "T25: Mixer detection", "T27: Cross-chain bridging detection"]
    },
    "DOMAIN": {
        "tasks": ["T28: WHOIS lookup", "T29: DNS enumeration", "T30: SSL certificate analysis", "T32: Hosting/IP correlation"]
    },
    "IP_ADDRESS": {
        "tasks": ["T34: Geolocation", "T35: ISP identification", "T36: Threat scoring", "T38: Infrastructure mapping"]
    },
    "EMAIL": {
        "tasks": ["T39: Breach lookup", "T40: Domain correlation", "T41: Alias detection", "T42: Social account linking"]
    },
    "PHONE": {
        "tasks": ["T44: Carrier lookup", "T46: Social linkage", "T48: Communication graph mapping"]
    },
    "FILE_DOC": {
        "tasks": ["T49: Metadata extraction", "T50: Hash generation", "T52: IoC extraction"]
    },
    "MALWARE_IOC": {
        "tasks": ["T64: Threat attribution", "T65: Signature matching", "T67: Campaign detection"]
    },
    "NETWORK_ENTITY": {
        "tasks": [
            "T71: Network topology mapping", 
            "T72: BGP & ASN route profiling",
            "T73: Port & protocol enumeration",
            "T74: Hidden Service (Tor/I2P) de-anonymization",
            "T75: Node proximity & lateral movement tracing"
        ]
    }
}

# Regex Type -> Ontology Class Mapping
ONTOLOGY_MAP = {
    "person": "PERSON", "usernames": "PERSON", "ssn": "PERSON", "device": "PERSON",
    "org": "ORGANIZATION",
    "domain": "DOMAIN", "url_path": "DOMAIN",
    "email": "EMAIL",
    "ip": "IP_ADDRESS",
    "eth": "CRYPTO_WALLET", "btc": "CRYPTO_WALLET", "sol": "CRYPTO_WALLET",
    "bitcoin": "CRYPTO_WALLET", "ethereum": "CRYPTO_WALLET", "bsc": "CRYPTO_WALLET", 
    "solana": "CRYPTO_WALLET", "ripple": "CRYPTO_WALLET", "avalanche": "CRYPTO_WALLET", 
    "tron": "CRYPTO_WALLET", "polygon": "CRYPTO_WALLET", "litecoin": "CRYPTO_WALLET", 
    "bitcoin_cash": "CRYPTO_WALLET", "stellar": "CRYPTO_WALLET", "bsv": "CRYPTO_WALLET", 
    "multiversx": "CRYPTO_WALLET", "algorand": "CRYPTO_WALLET", "flow": "CRYPTO_WALLET", 
    "klaytn": "CRYPTO_WALLET", "bitcoingold": "CRYPTO_WALLET", "zcash": "CRYPTO_WALLET", 
    "dash": "CRYPTO_WALLET", "harmony": "CRYPTO_WALLET", "verge": "CRYPTO_WALLET", 
    "provenance": "CRYPTO_WALLET", "sui": "CRYPTO_WALLET", "dogecoin": "CRYPTO_WALLET", 
    "cardano": "CRYPTO_WALLET", "conflux": "CRYPTO_WALLET", "ethereum_classic": "CRYPTO_WALLET", 
    "celo": "CRYPTO_WALLET", "fantom": "CRYPTO_WALLET", "moonbeam": "CRYPTO_WALLET", 
    "cronos": "CRYPTO_WALLET", "everscale": "CRYPTO_WALLET", "filecoin": "CRYPTO_WALLET", 
    "hedera": "CRYPTO_WALLET", "tezos": "CRYPTO_WALLET", "coreum": "CRYPTO_WALLET", "coti": "CRYPTO_WALLET",
    "phone": "PHONE",
    "documents": "FILE_DOC",
    "hash": "MALWARE_IOC",
    "onion": "DOMAIN",
    "monero": "CRYPTO_WALLET",
    "telegram_handle": "PERSON",
    "tox_ids": "NETWORK_ENTITY",
    "uuid": "NETWORK_ENTITY",
    "profile_id": "NETWORK_ENTITY",
    "credit_cards": "BANK_ACCOUNT"
}

# 🧠 UIE REGEX ENGINE (v1.2)
PATTERNS = {
    "person": r"\b(?!Inc|Ltd|LLC|Corp|Company|Department|University|Institute)([A-Z][a-z]+(?:\s[A-Z][a-z]+){1,3})\b",
    "org": r"\b([A-Z][A-Za-z0-9&.,'-]+\s(?:Inc|LLC|Ltd|Corp|Corporation|Foundation|Group|Agency|Institute|University))\b",
    "domain": r"\b((?:https?:\/\/)?(?:www\.)?[a-zA-Z0-9-]+\.[a-z]{2,}(?:\.[a-z]{2,})?)\b",
    "email": r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b",
    "ip": r"\b(?:(?:25[0-5]|2[0-4]\d|1?\d?\d)\.){3}(?:25[0-5]|2[0-4]\d|1?\d?\d)\b",
    "url_path": r"https?:\/\/(?:www\.)?[^\s/$.?#].[^\s]*",
    "bitcoin": r"\b(?:bc1[a-z0-9]{25,39}|[13][a-km-zA-HJ-NP-Z1-9]{25,34})\b",
    "ethereum": r"\b0x[a-fA-F0-9]{40}\b",
    "bsc": r"\b0x[a-fA-F0-9]{40}\b",
    "solana": r"\b[1-9A-HJ-NP-Za-km-z]{32,44}\b",
    "ripple": r"\br[1-9A-HJ-NP-Za-km-z]{24,34}\b",
    "avalanche": r"\b0x[a-fA-F0-9]{40}\b",
    "tron": r"\bT[1-9A-HJ-NP-Za-km-z]{33}\b",
    "polygon": r"\b0x[a-fA-F0-9]{40}\b",
    "litecoin": r"\b(?:ltc1[a-z0-9]{25,65}|[LM3][a-km-zA-HJ-NP-Z1-9]{26,33})\b",
    "bitcoin_cash": r"\b(?:bitcoincash:)?(?:q|p)[a-z0-9]{41}\b",
    "stellar": r"\bG[A-Z2-7]{55}\b",
    "bsv": r"\b[13][a-km-zA-HJ-NP-Z1-9]{25,34}\b",
    "multiversx": r"\berd1[a-z0-9]{58}\b",
    "algorand": r"\b[A-Z2-7]{58}\b",
    "flow": r"\b0x[a-fA-F0-9]{8,64}\b",
    "klaytn": r"\b0x[a-fA-F0-9]{40}\b",
    "bitcoingold": r"\b[13][a-km-zA-HJ-NP-Z1-9]{25,34}\b",
    "zcash": r"\b(?:t1|t3)[a-km-zA-HJ-NP-Z1-9]{25,34}\b",
    "dash": r"\b[Xx][a-km-zA-HJ-NP-Z1-9]{25,34}\b",
    "harmony": r"\b0x[a-fA-F0-9]{40}\b",
    "verge": r"\b[DdXxYyZz][a-km-zA-HJ-NP-Z1-9]{25,34}\b",
    "provenance": r"\b0x[a-fA-F0-9]{40}\b",
    "sui": r"\b0x[a-fA-F0-9]{40,64}\b",
    "dogecoin": r"\bD{1}[a-km-zA-HJ-NP-Z1-9]{25,34}\b",
    "cardano": r"\b(?:addr1|addr_test1)[0-9a-z]{58,}\b",
    "conflux": r"\b(?:cfx:)?0x[a-fA-F0-9]{40}\b",
    "ethereum_classic": r"\b0x[a-fA-F0-9]{40}\b",
    "celo": r"\b0x[a-fA-F0-9]{40}\b",
    "fantom": r"\b0x[a-fA-F0-9]{40}\b",
    "moonbeam": r"\b0x[a-fA-F0-9]{40}\b",
    "cronos": r"\b0x[a-fA-F0-9]{40}\b",
    "everscale": r"\b0:[a-fA-F0-9]{64,}\b",
    "filecoin": r"\b(?:f|t)[0-9][a-z0-9]{20,}\b",
    "hedera": r"\b0\.0\.\d+\b",
    "tezos": r"\b(?:tz1|tz2|tz3|KT1)[1-9A-HJ-NP-Za-km-z]{33}\b",
    "coreum": r"\b0x[a-fA-F0-9]{40}\b",
    "coti": r"\b[1-9A-HJ-NP-Za-km-z]{25,45}\b",
    "device": r"\b([A-Z0-9]{8,}-[A-Z0-9]{4,}-[A-Z0-9]{4,}-[A-Z0-9]{4,}-[A-Z0-9]{12,})\b",
    "event": r"\b([A-Z][a-zA-Z]+(?:\s[A-Z][a-zA-Z]+)*\s(?:202[0-9]|19[0-9]{2}))\b",
    "hash": r"\b(?<!0x)([a-fA-F0-9]{32}|[a-fA-F0-9]{40}|[a-fA-F0-9]{64})\b",
    "onion": r"\b[a-z2-7]{16,56}\.onion\b",
    "monero": r"\b(?:4[0-9AB][1-9A-HJ-NP-Za-km-z]{93}|8[1-9A-HJ-NP-Za-km-z]{94})\b",
    "telegram_handle": r"(?:https?:\/\/)?(?:t\.me|telegram\.me)\/([a-zA-Z0-9_]{5,32})\b",
    "tox_ids": r"\b[0-9A-Fa-f]{76}\b",
    "uuid": r"\b[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}\b",
    "profile_id": r"(?:(?<=-)|(?<=_))\d{10,16}\b",
    "credit_cards": {
        "visa": r"\b4[0-9]{12}(?:[0-9]{3})?\b",
        "mastercard": r"\b5[1-5][0-9]{14}\b",
        "amex": r"\b3[47][0-9]{13}\b",
        "discover": r"\b6(?:011|5[0-9]{2})[0-9]{12}\b",
    },
    "ssn": r"\b(?!000|666)[0-8][0-9]{2}-(?!00)[0-9]{2}-(?!0000)[0-9]{4}\b",
    "phone": r"\b(?:\+?1[-.\s]?)?(?:\(?\d{3}\)?[-.\s]?)?\d{3}[-.\s]?\d{4}\b",
    "documents": r"\b\S+\.(?:pdf|xml|docx?|txt|csv|xlsx?)\b",
    "usernames": r"(?<=^|(?<=[^a-zA-Z0-9-_\.]))@([A-Za-z0-9_]{3,40})"
}

def analyze_sentiment(text):
    """Analyze sentiment and threat intent of extracted text using Gemini LLM."""
    if not text: return "UNKNOWN"
    api_key = os.getenv("VITE_GEMINI_API_KEY", "")
    if not api_key: return "UNAVAILABLE (NO API KEY)"
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
    prompt = f"Analyze the overall sentiment and intent of the following darknet/OSINT text. Reply ONLY with one of the following exact labels: HIGHLY_MALICIOUS, SUSPICIOUS, NEUTRAL, POSITIVE. Text snippet: {text[:3000]}"
    try:
        r = requests.post(url, json={"contents": [{"parts": [{"text": prompt}]}]}, timeout=10)
        if r.status_code == 200:
            data = r.json()
            return data['candidates'][0]['content']['parts'][0]['text'].strip()
    except Exception as e:
        import logging, traceback
        logging.error(f'[Recovered Exception in darknetv2.py] {e}')
        traceback.print_exc()
    return "ANALYSIS_FAILED"

# Runtime structures
session = requests.Session()
if USE_TOR:
    session.proxies.update(TOR_PROXY)
session.headers.update(HEADERS)
try:
    ua = UserAgent()
except Exception:
    ua = None

queue_urls = Queue()
queued_set = set()
processed_set = set()

locks = {
    "processed": threading.Lock(),
    "queued": threading.Lock(),
    "state": threading.Lock(),
    "stats": threading.Lock(),
    "log": threading.Lock(),
    "sse": threading.Lock()
}

keywords = []
stats = {
    "seed": 0, "processed": 0, "pending": 0, 
    "btc_found": 0, "eth_found": 0, "sol_found": 0, "emails_found": 0, 
    "docs_found": 0, "orgs_found": 0, "hashes_found": 0
}

# Global event stream for UI
recent_logs = deque(maxlen=10)
recent_entities = deque(maxlen=15) # Holds tuples of (ts, ont_class, val, url)
graph_state_lock = threading.Lock()
graph_snapshot = {"nodes": [], "edges": []}
graph_index = {}
sse_clients = []

# Global executor for heavy OSINT background tasks
osint_background_pool = concurrent.futures.ThreadPoolExecutor(max_workers=10)

# Crawler Global State
global_worker_stop = threading.Event()
global_workers = []
crawler_is_running = False

def trigger_social_scan(val, node_id):
    def _scan():
        try:
            results = social_scanner.verify_username(val, max_workers=10)
            if results:
                with graph_state_lock:
                    existing = graph_index.get(node_id)
                    if existing:
                        existing.setdefault("metadata", {})["verified_socials"] = results
        except Exception:
            pass
    osint_background_pool.submit(_scan)

# -------------------- 🧠 UIE REGEX ENGINE & HEURISTICS --------------------
class UIEEngine:
    boilerplate_blacklist = {
        "market", "crypto", "bitcoin", "solana", "ethereum", "prediction", "trends", "world",
        "policy", "privacy", "terms", "service", "home", "search", "browse", "trending", 
        "holders", "positions", "activity", "finance", "geopolitics", "tech", "culture",
        "sports", "odds", "weather", "elections", "contract", "integrity", "price", "login",
        "sign", "account", "careers", "press", "adventure", "popular", "liquid", "ending",
        "soon", "about", "contact", "support", "faq", "frequently", "asked", "questions",
        "perps", "perpetual", "earnings", "economic", "news", "view", "more", "menu", "skip",
        "generate", "related", "above", "past", "may", "this", "what", "beware", "newest", "live", 
        "breaking", "fc", "stade", "rennais", "learn", "politics", "middle", "east", "polymarket",
        "esports", "iran", "rules", "context"
    }

    @staticmethod
    def validate_person(text):
        """Filters out UI navigation items and boilerplate flagged as Person."""
        clean_text = text.strip()
        if len(clean_text) < 3 or len(clean_text) > 35: return False
        words = clean_text.lower().split()
        if any(word in UIEEngine.boilerplate_blacklist for word in words): return False
        if "buy" in words or "sell" in words or "yes" in words or "no" in words: return False
        if clean_text.isupper(): return False
        return True

    @staticmethod
    def score_entity(e_type, value):
        confidence = 0.5
        if e_type in ['ip', 'hash', 'tox_ids', 'uuid', 'profile_id'] or ONTOLOGY_MAP.get(e_type) == 'CRYPTO_WALLET':
            confidence = 0.95
        if e_type == 'org' and any(suffix in value for suffix in ['Inc', 'LLC', 'Corp']):
            confidence += 0.3
        if e_type == 'person':
            words = value.split()
            if len(words) >= 2 and all(w[0].isupper() for w in words):
                confidence += 0.2
            else:
                confidence -= 0.2
        if e_type == 'email' and '@' in value:
            confidence = 0.9
        
        return round(min(1.0, max(0.1, confidence)), 2)

    @staticmethod
    def extract(text, source_url):
        import dataclasses
        import uuid
        
        uie_master = []
        ts = datetime.now(timezone.utc).isoformat()
        
        for e_type, pattern in PATTERNS.items():
            category = ONTOLOGY_MAP.get(e_type, "UNKNOWN")
            tasks = NEMESIS_ONTOLOGY.get(category, {}).get("tasks", [])
            
            if e_type == 'credit_cards':
                for cc_type, cc_pattern in pattern.items():
                    matches = list(set(re.findall(cc_pattern, text)))
                    for match in matches:
                        entity = IntelligenceEntity(
                            entity_id=str(uuid.uuid4()),
                            entity_type="BANK_ACCOUNT",
                            value=match,
                            confidence=0.9,
                            first_seen=ts,
                            last_seen=ts,
                            source_count=1,
                            evidence=[{"source_url": source_url, "timestamp": ts}],
                            attributes={"type": "credit_card", "subtype": cc_type, "autonomous_tasks": NEMESIS_ONTOLOGY.get("BANK_ACCOUNT", {}).get("tasks", [])}
                        )
                        e_dict = dataclasses.asdict(entity)
                        e_dict["type"] = "credit_card"
                        e_dict["ontology_class"] = "BANK_ACCOUNT"
                        e_dict["sourceSpan"] = source_url
                        e_dict["timestamp"] = ts
                        e_dict["autonomous_tasks"] = NEMESIS_ONTOLOGY.get("BANK_ACCOUNT", {}).get("tasks", [])
                        uie_master.append(e_dict)
            else:
                matches = list(set(re.findall(pattern, text)))
                for match in matches:
                    val = match[0] if isinstance(match, tuple) else match
                    val = val.strip()
                    if len(val) < 2: continue
                    
                    # HEURISTIC CONTEXT FILTERING
                    if e_type == "person":
                        if not UIEEngine.validate_person(val):
                            continue
                            
                    base_confidence = UIEEngine.score_entity(e_type, val)
                    
                    # Enhance confidence with Global Ontology
                    if category == "CRYPTO_WALLET":
                        base_confidence = max(base_confidence, CONFIDENCE.get("BLOCKCHAIN", 90) / 100.0)
                    elif category == "DOMAIN" or category == "IP_ADDRESS":
                        base_confidence = max(base_confidence, CONFIDENCE.get("OSINT", 50) / 100.0)
                            
                    entity = IntelligenceEntity(
                        entity_id=str(uuid.uuid4()),
                        entity_type=category,
                        value=val,
                        confidence=base_confidence,
                        first_seen=ts,
                        last_seen=ts,
                        source_count=1,
                        evidence=[{"source_url": source_url, "timestamp": ts}],
                        attributes={"type": e_type, "autonomous_tasks": tasks}
                    )
                    # We inject the legacy fields to ensure compatibility with graph engine
                    e_dict = dataclasses.asdict(entity)
                    e_dict["type"] = e_type
                    e_dict["ontology_class"] = category
                    e_dict["sourceSpan"] = source_url
                    e_dict["timestamp"] = ts
                    e_dict["autonomous_tasks"] = tasks
                    
                    uie_master.append(e_dict)
        
        filtered_master = [e for e in uie_master if e['confidence'] >= 0.4]
        return filtered_master

class SigmaGraphEngine:
    # Defines styles, icons, and hierarchical groupings for network graphing engines
    ONTOLOGY_STYLING = {
        "CRYPTO_WALLET": {"color": "#F7931A", "icon_name": "currency-btc", "icon_code": "\uf15a", "shape": "icon", "level": 4},
        "PERSON": {"color": "#00FFFF", "icon_name": "user", "icon_code": "\uf007", "shape": "icon", "level": 2},
        "EMAIL": {"color": "#FF00FF", "icon_name": "envelope", "icon_code": "\uf0e0", "shape": "icon", "level": 3},
        "BANK_ACCOUNT": {"color": "#00FF00", "icon_name": "building-columns", "icon_code": "\uf19c", "shape": "icon", "level": 4},
        "ORGANIZATION": {"color": "#FFFFFF", "icon_name": "building", "icon_code": "\uf1ad", "shape": "icon", "level": 1},
        "DOMAIN": {"color": "#3366FF", "icon_name": "globe", "icon_code": "\uf0ac", "shape": "icon", "level": 0},
        "IP_ADDRESS": {"color": "#FF4500", "icon_name": "network-wired", "icon_code": "\uf6ff", "shape": "icon", "level": 0},
        "FILE_DOC": {"color": "#808080", "icon_name": "file-lines", "icon_code": "\uf15c", "shape": "icon", "level": 5},
        "MALWARE_IOC": {"color": "#FF0000", "icon_name": "bug", "icon_code": "\uf088", "shape": "icon", "level": 5},
        "NETWORK_ENTITY": {"color": "#800080", "icon_name": "server", "icon_code": "\uf233", "shape": "icon", "level": 2},
        "PHONE": {"color": "#00FF7F", "icon_name": "phone", "icon_code": "\uf095", "shape": "icon", "level": 3},
        "UNKNOWN": {"color": "#555555", "icon_name": "circle-question", "icon_code": "\uf059", "shape": "dot", "level": 6}
    }

    @staticmethod
    def normalize_node_id(value, ont_class):
        return hashlib.sha256(f"{value}:{ont_class}".encode("utf-8")).hexdigest()

    @staticmethod
    def stream_nodes(uie_entities):
        if not uie_entities:
            return

        record_nodes = []
        for entity in uie_entities:
            ont_class = entity.get("ontology_class", "UNKNOWN")
            node_id = SigmaGraphEngine.normalize_node_id(entity.get("value", ""), ont_class)
            
            # Retrieve layout and rendering style config
            style = SigmaGraphEngine.ONTOLOGY_STYLING.get(ont_class, SigmaGraphEngine.ONTOLOGY_STYLING["UNKNOWN"])
            confidence_score = float(entity.get("confidence", 0.0))
            
            node = {
                "id": node_id,
                "label": str(entity.get("value", "")),
                "class": ont_class,
                "type": entity.get("type"),
                "confidence": confidence_score,
                "source": entity.get("sourceSpan"),
                "tasks": entity.get("autonomous_tasks", []),
                "group": ont_class,
                "color": style["color"],
                "shape": style["shape"],
                "level": style["level"], # Enables hierarchical structural layout
                "size": max(12, min(40, int(confidence_score * 30))), # Size scales with entity confidence
                "icon": {
                    "face": "FontAwesome", 
                    "code": style["icon_code"], 
                    "color": style["color"],
                    "name": style["icon_name"]
                },
                "metadata": {
                    k: v for k, v in entity.items() if k not in ("value", "ontology_class", "type", "confidence", "sourceSpan", "autonomous_tasks")
                }
            }
            record_nodes.append(node)

            with graph_state_lock:
                existing = graph_index.get(node_id)
                if existing:
                    existing["confidence"] = max(existing["confidence"], node["confidence"])
                    existing["tasks"] = list({*existing.get("tasks", []), *node.get("tasks", [])})
                    existing["size"] = max(12, min(40, int(existing["confidence"] * 30)))
                else:
                    graph_index[node_id] = node
                    graph_snapshot["nodes"].append(node)

        graph_update = {"nodes": [], "edges": []}
        seen_edge = set()
        for a in record_nodes:
            for b in record_nodes:
                if a["id"] == b["id"]:
                    continue
                pair = tuple(sorted([a["id"], b["id"]]))
                if pair in seen_edge:
                    continue
                seen_edge.add(pair)
                
                # Enhanced structural relationship mapping
                is_org_rel = a["class"] == "ORGANIZATION" or b["class"] == "ORGANIZATION"
                is_seq_rel = a["level"] != b["level"]
                
                if is_org_rel:
                    rel_type = "organization-link"
                    edge_color = "#FFFFFF"
                    edge_color = "#111111"
                elif is_seq_rel:
                    rel_type = "hierarchical-link"
                    edge_color = "#3399FF"
                    edge_color = "#0055CC"
                else:
                    rel_type = "co-occurrence"
                    edge_color = "#888888"
                    edge_color = "#666666"
                    
                edge = {
                    "source": pair[0], 
                    "target": pair[1], 
                    "type": rel_type, 
                    "count": 1,
                    "value": 1,
                    "dashes": a["class"] == "UNKNOWN" or b["class"] == "UNKNOWN",
                    "color": {"color": edge_color, "highlight": "#FF3333"},
                    "smooth": {"type": "cubicBezier"} if is_seq_rel else {"type": "continuous"}
                }
                
                with graph_state_lock:
                    existing_edge = next((e for e in graph_snapshot["edges"] if e["source"] == edge["source"] and e["target"] == edge["target"]), None)
                    if existing_edge:
                        existing_edge["count"] += 1
                        existing_edge["value"] = min(15, existing_edge["count"]) # Thicken line strictly on recurrent hits
                    else:
                        graph_snapshot["edges"].append(edge)
                graph_update["edges"].append(edge)

        graph_update["nodes"] = record_nodes
        if graph_update["nodes"] or graph_update["edges"]:
            push_sse_event("graph", graph_update)

# -------------------- 📄 DOCUMENT PARSING ENGINE --------------------
class DocumentParsingEngine:
    def __init__(self, download_dir="nemesis_vault", max_file_size_mb=25):
        self.download_dir = download_dir
        self.max_file_size = max_file_size_mb * 1024 * 1024
        os.makedirs(self.download_dir, exist_ok=True)

    def get_extension_from_url(self, url):
        """Extracts extension or infers type from the URL path."""
        path = urlparse(url).path
        ext = os.path.splitext(path)[1].lower()
        return ext

    def parse_binary_stream(self, file_bytes, extension):
        """Parses binary data in-memory based on file extension and returns text."""
        extracted_text = ""
        try:
            if extension == ".pdf":
                pdf_reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                text_layers = []
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_layers.append(page_text)
                extracted_text = "\n".join(text_layers)

            elif extension in [".docx", ".doc"]:
                doc = docx.Document(io.BytesIO(file_bytes))
                extracted_text = "\n".join([para.text for para in doc.paragraphs])

            elif extension in [".pptx", ".ppt"]:
                prs = pptx.Presentation(io.BytesIO(file_bytes))
                text_layers = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_layers.append(shape.text)
                extracted_text = "\n".join(text_layers)

            elif extension == ".csv":
                decoded_csv = file_bytes.decode('utf-8', errors='ignore')
                csv_reader = csv.reader(io.StringIO(decoded_csv))
                extracted_text = "\n".join([",".join(row) for row in csv_reader])

            elif extension in [".xlsx", ".xls"]:
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
                text_layers = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join([str(cell) for cell in row if cell is not None])
                        if row_text:
                            text_layers.append(row_text)
                extracted_text = "\n".join(text_layers)

            elif extension in [".txt", ".log", ".ini", ".conf"]:
                extracted_text = file_bytes.decode('utf-8', errors='ignore')

        except Exception as e:
            return f"[PARSING_ERROR] Failed to parse {extension} file: {str(e)}"
        
        return extracted_text

# Initialize globally
doc_parser = DocumentParsingEngine(download_dir="nemesis_vault", max_file_size_mb=20)


# -------------------- PARALLEL API ENRICHMENT ENGINE --------------------
class WalletEnrichmentEngine:
    def __init__(self):
        self.providers = {
            "ETH": [
                {"name": "Etherscan", "url": "https://api.etherscan.io/api?module=account&action=txlist&address={}&apikey={key}", "key": os.getenv("VITE_ETHERSCAN_API_KEY", "")},
                {"name": "Blockscout", "url": "https://eth.blockscout.com/api?module=account&action=txlist&address={}", "key": ""}
            ],
            "BTC": [
                {"name": "Blockchain.info", "url": "https://blockchain.info/rawaddr/{}", "key": ""},
                {"name": "BlockCypher", "url": "https://api.blockcypher.com/v1/btc/main/addrs/{}", "key": ""}
            ]
        }
        self.executor = concurrent.futures.ThreadPoolExecutor(max_workers=5)

    def fetch_data(self, address, network):
        providers = self.providers.get(network.upper(), [])
        if not providers:
            return {"status": "skipped", "message": f"No providers for {network}"}

        for provider in providers:
            try:
                url = provider["url"].format(address, key=provider["key"]) if provider["key"] else provider["url"].format(address)
                headers = {"User-Agent": "Nemesis-OSINT-Engine/3.0"}
                resp = requests.get(url, headers=headers, timeout=5)
                
                if resp.status_code == 200:
                    data = resp.json()
                    tx_count = 0
                    balance = "0"
                    
                    if provider["name"] in ["Etherscan", "Blockscout"]:
                        txs = data.get("result", [])
                        tx_count = len(txs) if isinstance(txs, list) else 0
                    elif provider["name"] == "Blockchain.info":
                        tx_count = data.get("n_tx", 0)
                        balance = str(data.get("final_balance", 0))
                    
                    # --- AUTONOMOUS TRACING & CEX IDENTIFICATION ANALYSIS ---
                    addr_lower = address.lower()
                    identified_cex = KNOWN_ENTITIES.get(addr_lower)
                    if not identified_cex:
                        if tx_count > 10000:
                            identified_cex = "Large Custodial Hub"
                        elif tx_count > 500:
                            identified_cex = "Active EOA / Pool Node"
                        else:
                            identified_cex = "Private Wallet Node"
                            
                    cex_hash = hashlib.md5(identified_cex.encode()).hexdigest()[:10]
                    total_landed = "0.00"
                    if provider["name"] == "Blockchain.info" and balance != "0":
                        try:
                            total_landed = f"{int(balance) / 1e8:.4f}"
                        except Exception as e:
                            import logging, traceback
                            logging.error(f'[Recovered Exception in darknetv2.py] {e}')
                            traceback.print_exc()
                    
                    tx_graph = None
                    if tx_count > 0:
                        tx_graph = {
                            "nodes": [
                                {"id": address, "label": f"Target\n{address[:6]}...", "color": "#F7931A", "shape": "box", "font": {"color": "#000"}},
                                {"id": f"node_{cex_hash}", "label": f"{identified_cex}\n(Identified Endpoint)", "color": "#32CD32", "shape": "box", "font": {"color": "#000", "bold": True}}
                            ],
                            "edges": [
                                {"source": address, "target": f"node_{cex_hash}", "label": f"{total_landed} {network}", "arrows": "to", "color": "#FF4500"}
                            ]
                        }

                    return {
                        "status": "success",
                        "provider": provider["name"],
                        "tx_count": tx_count,
                        "balance": balance,
                        "identified_cex": identified_cex,
                        "total_landed_amount": f"{total_landed} {network}",
                        "tx_graph": tx_graph
                    }
                elif resp.status_code == 429:
                    continue # Rate limited, instantly failover to fallback provider
            except Exception:
                continue # Connection error, instantly failover to fallback provider
        return {"status": "error", "message": "All API providers failed or rate-limited"}

    def enrich(self, entities):
        wallets = [e for e in entities if e.get("ontology_class") == "CRYPTO_WALLET"]
        if not wallets:
            return entities
            
        future_to_wallet = {}
        for w in wallets:
            addr = w["value"]
            net = "ETH" if str(addr).startswith("0x") else "BTC"
            future = self.executor.submit(self.fetch_data, addr, net)
            future_to_wallet[future] = w

        for future in concurrent.futures.as_completed(future_to_wallet):
            w = future_to_wallet[future]
            try:
                res = future.result()
                if res["status"] == "success":
                    w["enrichment"] = res
                    # Auto-label based on on-chain activity volume
                    if res.get("tx_count", 0) > 50:
                        w["tags"] = ["HIGH_ACTIVITY_NODE"]
                    elif res.get("tx_count", 0) == 0:
                        w["tags"] = ["INACTIVE"]
                    else:
                        w["tags"] = ["ACTIVE"]
            except Exception:
                pass
        
        return entities

enrichment_engine = WalletEnrichmentEngine()

# -------------------- OSINT MAPPER SUBPROCESS (AUTO-RECON) --------------------
def run_osint_mapper():
    """Runs exclusively in the secondary console window to display OSINT Recon results."""
    console.print(Panel.fit("[bold magenta]🧠 NEMESIS AUTONOMOUS OSINT MAPPER (v3.0)[/bold magenta]"))
    console.print("Tailing `export.jsonl` for UIE Identity Cross-Referencing & Auto-Recon...\n")
    
    platforms = {
        "GitHub": "https://github.com/{}",
        "Twitter": "https://twitter.com/{}",
        "Reddit": "https://www.reddit.com/user/{}/"
    }

    # Wait until export.jsonl is created by the main process
    while not os.path.exists(EXPORT_JSONL):
        time.sleep(1)

    # Use a standard tail -f reading mechanism to prevent lagging/re-reading the whole file
    with open(EXPORT_JSONL, "r", encoding="utf-8") as f:
        while True:
            line = f.readline()
            if not line:
                time.sleep(0.5)
                continue
            
            line = line.strip()
            if not line:
                continue
                
            try:
                record = json.loads(line)
                web_info = record.get("web_info", {})
                metadata = web_info.get("metadata", {})
                url = web_info.get("url", "Unknown")
                uie_entities = record.get("uie_entities", [])
                
                # If no entities and no metadata, skip to avoid spamming the console
                if not uie_entities and not metadata:
                    continue
                
                # --- 1. Display Target and Scraped Metadata ---
                console.print(Panel(f"[bold green]Target Node:[/bold green] {url}", border_style="green"))
                
                if metadata:
                    meta_tree = Tree("[bold cyan]📄 Page Meta-Data Extracted[/bold cyan]")
                    for k, v in metadata.items():
                        meta_tree.add(f"[dim]{k}:[/dim] {v}")
                    console.print(meta_tree)
                
                # --- 2. Auto-Recon on Entities ---
                for entity in uie_entities:
                    if entity.get('confidence', 0) < 0.8: continue
                    
                    val = entity.get('value', 'Unknown')
                    ont_class = entity.get('ontology_class', 'UNKNOWN')
                    tasks = entity.get('autonomous_tasks', [])
                    
                    tree = Tree(f"[bold cyan]Entity Discovered:[/bold cyan] {val} [yellow]({ont_class})[/yellow]")
                    
                    # Display auto-labels and parallel enrichment data if available
                    if 'enrichment' in entity:
                        enrich_data = entity['enrichment']
                        tags = entity.get('tags', [])
                        tag_str = f" [bold red]{tags}[/bold red]" if tags else ""
                        tree.add(f"[bold green]📊 On-Chain Stats:[/bold green] {enrich_data.get('tx_count')} TXs via {enrich_data.get('provider')}{tag_str}")

                    recon_branch = tree.add("[bold yellow]🔍 Auto-OSINT Recon Links & Actions:[/bold yellow]")
                    
                    # Standard Platform Lookups
                    if ont_class in ['PERSON', 'EMAIL', 'NETWORK_ENTITY']:
                        target = val.replace("@", "")
                        for plat_name, tmpl in platforms.items():
                            try:
                                p_url = tmpl.format(target) if "{}" in tmpl else tmpl + target
                                r = requests.get(p_url, headers=HEADERS, timeout=3)
                                if r.status_code == 200:
                                    recon_branch.add(f"[green]✔ {plat_name} Match![/green] -> {p_url}")
                            except Exception: pass
                            
                    # Specialized Passive Reconnaissance
                    if ont_class == 'IP_ADDRESS':
                        recon_branch.add(f"AbuseIPDB (Threat check): https://www.abuseipdb.com/check/{val}")
                        recon_branch.add(f"Shodan (Port check): https://www.shodan.io/host/{val}")
                    
                    elif ont_class == 'CRYPTO_WALLET':
                        recon_branch.add(f"Blockchair Explorer: https://blockchair.com/search?q={val}")
                        recon_branch.add(f"OXT Analytics: https://oxt.me/address/{val}")
                        
                    elif ont_class == 'MALWARE_IOC':
                        recon_branch.add(f"VirusTotal Analysis: https://www.virustotal.com/gui/search/{val}")
                        
                    elif ont_class == 'DOMAIN':
                        clean_domain = val.replace("https://", "").replace("http://", "").split("/")[0]
                        try:
                            resolved_ip = socket.gethostbyname(clean_domain)
                            recon_branch.add(f"[bold green]Resolved IP:[/bold green] {resolved_ip}")
                            recon_branch.add(f"Shodan: https://www.shodan.io/host/{resolved_ip}")
                        except Exception:
                            recon_branch.add("[dim]DNS Resolution Failed (Likely Hidden Service/Tor)[/dim]")
                        recon_branch.add(f"WHOIS Lookup: https://who.is/whois/{clean_domain}")
                        
                    elif ont_class == 'EMAIL':
                        recon_branch.add(f"HaveIBeenPwned: https://haveibeenpwned.com/account/{val}")

                    if tasks:
                        task_branch = tree.add("[bold magenta]⚡ Autonomous Task Pipeline:[/bold magenta]")
                        for task in tasks[:3]:
                            task_branch.add(f"[blue]Queued ->[/blue] {task}")
                            
                    console.print(tree)
                    console.print("-" * 70)
                    
            except json.JSONDecodeError: 
                pass
            except Exception as e:
                console.print(f"[red]Error processing mapped line: {e}[/red]")

def spawn_osint_mapper():
    script = os.path.abspath(sys.argv[0])
    try:
        if os.name == 'nt':
            # Uses subprocess.CREATE_NEW_CONSOLE to safely bypass cmd.exe quoting restrictions
            subprocess.Popen([sys.executable, script, "--osint-mapper"], creationflags=subprocess.CREATE_NEW_CONSOLE)
        elif sys.platform == 'darwin':
            cmd_args = f'"{sys.executable}" "{script}" --osint-mapper'
            subprocess.Popen(f'osascript -e \'tell app "Terminal" to do script "{cmd_args}"\'', shell=True)
        else:
            # Linux fallback list
            terminals = ['gnome-terminal', 'xfce4-terminal', 'konsole', 'x-terminal-emulator']
            launched = False
            for t in terminals:
                if os.system(f"which {t} > /dev/null") == 0:
                    subprocess.Popen([t, '-e', f'{sys.executable} {script} --osint-mapper'])
                    launched = True
                    break
            if not launched:
                log("[WARN] Could not find a compatible terminal to auto-spawn on Linux.", style="yellow")
                
        log("[SYSTEM] Auto-Recon OSINT Mapper console launched.", style="bold magenta")
    except Exception as e:
        log(f"[WARN] Could not auto-spawn mapper window: {e}", style="red")


# -------------------- FRONTEND SSE INTEGRATION --------------------
def push_sse_event(event_type, data):
    """Pushes real-time events to all connected React dashboards"""
    msg = f"event: {event_type}\ndata: {json.dumps(data)}\n\n"
    with locks["sse"]:
        for client_queue in sse_clients:
            client_queue.put(msg)

def log(msg, style="white"):
    ts = datetime.now().strftime("%H:%M:%S")
    with locks["log"]:
        recent_logs.append(f"[{ts}] [{style}]{msg}[/{style}]")
    push_sse_event("log", {"ts": ts, "message": msg, "style": style})

def record_entity(ont_class, value, source_url):
    ts = datetime.now().strftime("%H:%M:%S")
    with locks["log"]:
        display_url = source_url[:45] + "..." if len(source_url) > 45 else source_url
        recent_entities.append((ts, ont_class, value, display_url))
    push_sse_event("entity", {"ts": ts, "type": ont_class, "val": value, "src": source_url})

# -------------------- FLASK API SERVER --------------------
import logging
# Suppress default Flask/Werkzeug startup logs to prevent them from breaking the Rich UI
werkzeug_logger = logging.getLogger('werkzeug')
werkzeug_logger.setLevel(logging.ERROR)

# Safely suppress the Flask server banner without crashing Werkzeug environments
flask.cli.show_server_banner = lambda *args: None

# -------------------- AUTONOMOUS TASK EXECUTOR --------------------
class AutonomousTaskExecutor:
    def __init__(self, db_handler):
        self.dbh = db_handler
        self.task_log = deque(maxlen=100)
        self.task_lock = threading.Lock()

    def execute_tasks_for_entity(self, entity, source_url):
        """Execute autonomous tasks mapped to entity ontology class."""
        tasks = entity.get('autonomous_tasks', [])
        if not tasks:
            return

        entity_id = SigmaGraphEngine.normalize_node_id(entity.get('value', ''), entity.get('ontology_class', ''))
        for task in tasks:
            try:
                result = self._run_task(task, entity, source_url)
                with self.task_lock:
                    self.task_log.append({
                        "ts": datetime.now(timezone.utc).isoformat(),
                        "entity_id": entity_id,
                        "entity_val": entity.get('value'),
                        "task": task,
                        "status": "success" if result else "failed",
                        "result": result
                    })
                push_sse_event("task", {"entity": entity.get('value'), "task": task, "status": "success"})
            except Exception as e:
                log(f"[TASK] {task} failed for {entity.get('value')}: {str(e)[:60]}", style="dim red")

    def _run_task(self, task_name, entity, source_url):
        """Route to task-specific handler."""
        ont_class = entity.get('ontology_class', '')
        
        if ont_class == 'CRYPTO_WALLET':
            return self._verify_blockchain_address(entity.get('value'), task_name)
        elif ont_class == 'EMAIL':
            return self._verify_email(entity.get('value'), task_name)
        elif ont_class == 'DOMAIN':
            return self._verify_domain(entity.get('value'), task_name)
        elif ont_class == 'PERSON':
            return self._search_person_osint(entity.get('value'), task_name)
        elif ont_class in ['ORGANIZATION', 'NETWORK_ENTITY']:
            return self._lookup_organization(entity.get('value'), task_name)
        return None

    def _verify_blockchain_address(self, addr, task):
        if 'lookup' in task.lower() or 'verify' in task.lower():
            try:
                headers = {'User-Agent': ua.random if ua else 'Mozilla/5.0'}
                r = requests.get(f"https://api.blockchair.com/bitcoin/addresses/{addr}", timeout=5)
                return r.status_code == 200 if r else False
            except: return False
        return True

    def _verify_email(self, email, task):
        if 'verify' in task.lower():
            return bool(re.match(r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$', email))
        return True

    def _verify_domain(self, domain, task):
        if 'lookup' in task.lower() or 'whois' in task.lower():
            try:
                r = requests.get(f"https://dns.google/resolve?name={domain}", timeout=5)
                return r.status_code == 200 if r else False
            except: return False
        return True

    def _search_person_osint(self, name, task):
        """Implement the GLOBAL INTELLIGENCE SEARCH QUERY FRAMEWORK for Persons."""
        if not name: return False
        try:
            enc_name = urlparse(name).path if "://" in name else name.replace(" ", "+")
            
            # Generate OSINT Query URLs for Deep Scraping
            queries = [
                f"https://api.opencorporates.com/v0.4/officers/search?q={enc_name}",
                f"https://www.courtlistener.com/?q={enc_name}",
                f"https://html.duckduckgo.com/html/?q={enc_name}+SEC+indictment",
                f"https://sanctionssearch.ofac.treas.gov/Details.aspx?id={enc_name}"
            ]
            
            for q in queries:
                if q not in queued_set:
                    queued_set.add(q)
                    queue_urls.put(q)
            return True
        except: return False

    def _lookup_organization(self, org, task):
        """Implement the GLOBAL INTELLIGENCE SEARCH QUERY FRAMEWORK for Corporations."""
        if not org: return False
        try:
            enc_org = urllib.parse.quote(org)
            
            recon_queries = [
                f"https://html.duckduckgo.com/html/?q=site:sec.gov+%22{enc_org}%22",
                f"https://html.duckduckgo.com/html/?q=ext:pdf+%22{enc_org}%22+contract"
            ]
            
            for q in recon_queries:
                if q not in queued_set:
                    queued_set.add(q)
                    queue_urls.put(q)
            return True
        except: return False

    def get_task_log(self):
        with self.task_lock:
            return list(self.task_log)

# -------------------- GOOGLE SHEETS API INTEGRATION --------------------
class GoogleSheetsManager:
    def __init__(self):
        pass
    def authenticate(self):
        pass
    def export_entity(self, entity, url):
        pass
    def force_flush(self):
        return 0
    def clear_sheet(self):
        return False
    def sync_from_sheets(self, db_handler):
        pass

# -------------------- DOCUMENT EXTRACTION & STORAGE MANAGER --------------------
class DocumentStorageManager:
    def __init__(self, base_dir="nemesis_documents"):
        self.base_dir = base_dir
        os.makedirs(base_dir, exist_ok=True)
        self.storage_log = deque(maxlen=100)
        self.storage_lock = threading.Lock()

    def store_document(self, entity_id, entity_value, content, file_type, source_url):
        """Store extracted document with entity linkage."""
        try:
            entity_dir = os.path.join(self.base_dir, entity_id[:12])
            os.makedirs(entity_dir, exist_ok=True)
            
            filename = f"{entity_value.replace(' ', '_')[:30]}_{int(time.time())}.{file_type}"
            filepath = os.path.join(entity_dir, filename)
            
            if isinstance(content, str):
                with open(filepath, 'w', encoding='utf-8') as f:
                    f.write(content)
            else:
                with open(filepath, 'wb') as f:
                    f.write(content)
            
            with self.storage_lock:
                self.storage_log.append({
                    "ts": datetime.now(timezone.utc).isoformat(),
                    "entity_id": entity_id,
                    "entity": entity_value,
                    "file": filepath,
                    "type": file_type,
                    "source": source_url
                })
            
            log(f"[DOC] Extracted {file_type} → {filepath}", style="cyan")
            return filepath
        except Exception as e:
            log(f"[DOC] Store failed: {str(e)}", style="red")
            return None

    def get_storage_log(self):
        with self.storage_lock:
            return list(self.storage_log)

# -------------------- INTELLIGENCE PROFILE AGGREGATOR --------------------
class IntelligenceProfileAggregator:
    def __init__(self, dbh):
        self.dbh = dbh
        self.profiles = {}
        self.profile_lock = threading.Lock()

    def build_entity_profile(self, entity_id, entity):
        """Build complete intelligence dossier for an entity."""
        with self.profile_lock:
            profile = {
                "entity_id": entity_id,
                "value": entity.get('value'),
                "class": entity.get('ontology_class'),
                "confidence": entity.get('confidence', 0.0),
                "first_seen": datetime.now(timezone.utc).isoformat(),
                "sources": [entity.get('sourceSpan', '')],
                "relationships": [],
                "tasks_executed": [],
                "documents_extracted": [],
                "risk_score": 0.5
            }
            self.profiles[entity_id] = profile
            return profile

    def add_relationship(self, source_id, target_id, rel_type="co-occurrence"):
        with self.profile_lock:
            if source_id in self.profiles:
                self.profiles[source_id]["relationships"].append({
                    "target_id": target_id,
                    "type": rel_type
                })

    def add_task_result(self, entity_id, task, status, result):
        with self.profile_lock:
            if entity_id in self.profiles:
                self.profiles[entity_id]["tasks_executed"].append({
                    "task": task,
                    "status": status,
                    "result": result,
                    "timestamp": datetime.now(timezone.utc).isoformat()
                })

    def add_document(self, entity_id, doc_path):
        with self.profile_lock:
            if entity_id in self.profiles:
                self.profiles[entity_id]["documents_extracted"].append(doc_path)

    def get_profile(self, entity_id):
        with self.profile_lock:
            return self.profiles.get(entity_id)

    def get_all_profiles(self):
        with self.profile_lock:
            return list(self.profiles.values())

app = Flask(__name__)
CORS(app)

global_dbh = None
autonomy_executor = None
doc_storage = None
intel_aggregator = None
sheets_manager = None

@app.route('/')
def dashboard():
    return render_template('dashboard.html')

@app.route('/intelligence')
def intelligence_dashboard():
    try:
        return render_template('intelligence.html')
    except Exception as e:
        return f"<h2>Intelligence Dashboard Not Found</h2>", 404

@app.route('/api/stats')
def api_stats():
    with locks["stats"]:
        return jsonify({"status": "success", "metrics": stats})

@app.route('/api/vasp/search')
def api_vasp_search():
    query = request.args.get('q', '').strip()
    if not query or not global_dbh: return jsonify({"results": []})
    
    output = []
    try:
        if global_dbh.dbtype == "mongodb" and global_dbh.coll is not None:
            # Search in vasp, vasp_dir, coinbase_wallets
            collections_to_search = ["vasp", "vasp_dir", "coinbase_wallets", "entity", "wallet_labels"]
            mongo_query = {
                "$or": [
                    {"name": {"$regex": query, "$options": "i"}},
                    {"id": {"$regex": query, "$options": "i"}},
                    {"website": {"$regex": query, "$options": "i"}},
                    {"twitter": {"$regex": query, "$options": "i"}},
                    {"domain": {"$regex": query, "$options": "i"}},
                    {"entity": {"$regex": query, "$options": "i"}},
                    {"vasp": {"$regex": query, "$options": "i"}}
                ]
            }
            
            all_cols = global_dbh.db.list_collection_names()
            for col_name in collections_to_search:
                if col_name in all_cols:
                    try:
                        cursor = global_dbh.db[col_name].find(mongo_query).limit(50)
                        for doc in cursor:
                            doc['_id'] = str(doc.get('_id', ''))
                            output.append(doc)
                    except Exception: pass
    except Exception as e:
        print(f"MongoDB vasp search error: {e}")
    
    return jsonify({"results": output})

@app.route('/api/wallet/search')
def api_wallet_search():
    query = request.args.get('q', '').strip()
    if not query or not global_dbh: return jsonify({"results": []})
    
    output = []
    try:
        if global_dbh.dbtype == "mongodb" and global_dbh.coll is not None:
            # Search in entity, wallet_labels, coinbase_wallets
            collections_to_search = ["entity", "wallet_labels", "coinbase_wallets"]
            mongo_query = {
                "$or": [
                    {"address": {"$regex": query, "$options": "i"}},
                    {"label": {"$regex": query, "$options": "i"}},
                    {"name": {"$regex": query, "$options": "i"}},
                    {"category": {"$regex": query, "$options": "i"}}
                ]
            }
            
            all_cols = global_dbh.db.list_collection_names()
            for col_name in collections_to_search:
                if col_name in all_cols:
                    try:
                        cursor = global_dbh.db[col_name].find(mongo_query).limit(50)
                        for doc in cursor:
                            doc['_id'] = str(doc.get('_id', ''))
                            output.append(doc)
                    except Exception: pass
    except Exception as e:
        print(f"MongoDB wallet search error: {e}")
        
    return jsonify({"results": output})

@app.route('/api/search')
def api_search():
    query = request.args.get('q', '').strip()
    filter_type = request.args.get('filter', 'ALL').upper()
    try:
        page = int(request.args.get('page', 1))
        limit = int(request.args.get('limit', 20))
    except ValueError:
        page, limit = 1, 20
        
    skip_val = (page - 1) * limit
    
    if not query or not global_dbh:
        return jsonify({"results": [], "total": 0, "page": page, "limit": limit})
        
    try:
        output = []
        total_matches = 0
        if global_dbh.dbtype == "mongodb" and global_dbh.coll is not None:
            mongo_query = {
                "$or": [
                    {"content": {"$regex": query, "$options": "i"}},
                    {"title": {"$regex": query, "$options": "i"}},
                    {"description": {"$regex": query, "$options": "i"}},
                    {"url": {"$regex": query, "$options": "i"}},
                    {"web_info.content": {"$regex": query, "$options": "i"}},
                    {"web_info.title": {"$regex": query, "$options": "i"}},
                    {"web_info.description": {"$regex": query, "$options": "i"}},
                    {"web_info.url": {"$regex": query, "$options": "i"}},
                    {"keywords_detected": {"$regex": query, "$options": "i"}},
                    {"uie_entities.value": {"$regex": query, "$options": "i"}},
                    {"name": {"$regex": query, "$options": "i"}},
                    {"id": {"$regex": query, "$options": "i"}},
                    {"website": {"$regex": query, "$options": "i"}},
                    {"twitter": {"$regex": query, "$options": "i"}}
                ]
            }
            if filter_type != 'ALL':
                filter_map = {
                    'CRYPTO_WALLET': ['CRYPTO_WALLET', 'BANK_ACCOUNT', 'CRYPTO_TRANSACTION'],
                    'EMAIL/PERSON': ['PERSON', 'EMAIL'],
                    'MALWARE_IOC': ['MALWARE_IOC'],
                    'ORGANIZATION': ['ORGANIZATION']
                }
                allowed = filter_map.get(filter_type, [])
                if allowed:
                    mongo_query["uie_entities.ontology_class"] = {"$in": allowed}
            
            try:
                all_cols = global_dbh.db.list_collection_names()
                if 'darknet_data' in all_cols:
                    all_cols.remove('darknet_data')
                    all_cols.insert(0, 'darknet_data')
                
                for col_name in all_cols:
                    if len(output) >= limit:
                        break
                    try:
                        cursor = list(global_dbh.db[col_name].find(mongo_query).sort("crawled_at", -1).limit(limit - len(output)))
                        output.extend(cursor)
                    except Exception:
                        pass
                
                total_matches = len(output)
            except Exception as e:
                print(f"MongoDB search error: {e}")
                pass
            
            for doc in output:
                doc['_id'] = str(doc.get('_id', ''))
                if 'web_info' not in doc:
                    doc['web_info'] = {'url': doc.get('url', ''), 'title': doc.get('title', 'Untitled'), 'content': str(doc)[:200]}
            
            # Sort combined and apply skip/limit
            output.sort(key=lambda x: x.get('crawled_at', ''), reverse=True)
            output = output[skip_val:skip_val+limit]

        elif global_dbh.dbtype in ["sqlite", "postgresql"]:
            q = f"%{query}%"
            with global_dbh.lock:
                if global_dbh.dbtype == "sqlite":
                    global_dbh.cursor.execute("SELECT hash_id,crawled_at,url,title,description,content,links,financial,person,documents,keywords_detected FROM crawled WHERE content LIKE ? OR title LIKE ? OR description LIKE ? OR url LIKE ? LIMIT 50", (q,q,q,q))
                else:
                    global_dbh.cursor.execute("SELECT hash_id,crawled_at,url,title,description,content,links,financial,person,documents,keywords_detected FROM crawled WHERE content ILIKE %s OR title ILIKE %s OR description ILIKE %s OR url ILIKE %s LIMIT 50", (q,q,q,q))
                rows = global_dbh.cursor.fetchall()
            for row in rows:
                output.append({
                    "hash-ID": row[0],
                    "crawled_at": row[1],
                    "web_info": {"url": row[2], "title": row[3], "description": row[4], "content": row[5][:500]},
                    "uie_entities": json.loads(row[7] or "[]") + json.loads(row[8] or "[]") + json.loads(row[9] or "[]"),
                    "keywords_detected": json.loads(row[10] or "[]")
                })
        return jsonify({"results": output})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)}), 500

@app.route('/api/export_pdf')
def api_export_pdf():
    query = request.args.get('q', '').strip()
    if not query or not global_dbh:
        return "<h2>Query required for PDF Dossier compilation.</h2>", 400
        
    output = []
    try:
        # Standardized DB extraction query similar to normal search
        if global_dbh.dbtype == "mongodb" and global_dbh.coll is not None:
            search_pattern = re.compile(re.escape(query), re.IGNORECASE)
            q_filter = {"$or": [{"web_info.content": search_pattern}, {"web_info.title": search_pattern}, {"web_info.url": search_pattern}, {"uie_entities.value": search_pattern}, {"url": search_pattern}, {"title": search_pattern}]}
            cursor1 = list(global_dbh.db["darknet_data"].find(q_filter).sort("crawled_at", -1).limit(50))
            cursor2 = list(global_dbh.db["darknet_url"].find(q_filter).sort("crawled_at", -1).limit(50))
            for doc in cursor1 + cursor2:
                if 'web_info' not in doc:
                    doc['web_info'] = {'url': doc.get('url', ''), 'title': doc.get('title', 'Untitled'), 'content': str(doc)}
                output.append(doc)
            
        elif global_dbh.dbtype in ["sqlite", "postgresql"]:
            q = f"%{query}%"
            with global_dbh.lock:
                if global_dbh.dbtype == "sqlite":
                    global_dbh.cursor.execute("SELECT hash_id,crawled_at,url,title,description,content,links,financial,person,documents,keywords_detected FROM crawled WHERE content LIKE ? OR title LIKE ? OR description LIKE ? OR url LIKE ? LIMIT 50", (q,q,q,q))
                else:
                    global_dbh.cursor.execute("SELECT hash_id,crawled_at,url,title,description,content,links,financial,person,documents,keywords_detected FROM crawled WHERE content ILIKE %s OR title ILIKE %s OR description ILIKE %s OR url ILIKE %s LIMIT 50", (q,q,q,q))
                rows = global_dbh.cursor.fetchall()
            for row in rows:
                output.append({
                    "hash-ID": row[0],
                    "crawled_at": row[1],
                    "web_info": {"url": row[2], "title": row[3], "description": row[4], "content": row[5]},
                    "uie_entities": json.loads(row[7] or "[]") + json.loads(row[8] or "[]") + json.loads(row[9] or "[]"),
                    "keywords_detected": json.loads(row[10] or "[]")
                })
            
        html = f"""
        <!DOCTYPE html><html><head><meta charset="UTF-8"><title>NEMESIS OSINT DOSSIER: {query}</title>
        <style>
            body {{ font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif; padding: 40px; color: #333; }}
            h1 {{ color: #1e293b; border-bottom: 2px solid #cbd5e1; padding-bottom: 10px; text-transform: uppercase; }}
            .dossier {{ border: 1px solid #e2e8f0; padding: 20px; margin-bottom: 20px; page-break-inside: avoid; border-radius: 8px; background: #f8fafc; }}
            .meta {{ font-size: 12px; color: #64748b; margin-bottom: 15px; }}
            .content {{ font-size: 14px; background: #fff; padding: 15px; border: 1px solid #e2e8f0; border-radius: 4px; }}
            .entities {{ margin-top: 15px; font-size: 12px; font-weight: bold; color: #b91c1c; }}
        </style></head><body onload="window.print()">
        <h1>Intel Dossier Report: {query}</h1>
        <p class="meta">Generated by Lionsgate NEMESIS OSINT Matrix | Date: {{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}}</p>
        """
        
        if not output:
            html += "<p>No threat intelligence dossiers found for the specified query.</p>"
            
        for doc in output:
            title = doc.get("web_info", {}).get("title", "Untitled")
            url = doc.get("web_info", {}).get("url", "Unknown")
            content = doc.get("web_info", {}).get("content", "")[:500]
            entities = [e.get("value") for e in doc.get("uie_entities", []) if e.get("confidence", 0) > 0.6]
            html += f"<div class='dossier'><h2>{{title}}</h2><p class='meta'>Source: {{url}}</p><div class='content'>{{content}}...</div>"
            if entities: html += f"<div class='entities'>High-Confidence Entities Mapped: {{', '.join(set(entities))}}</div>"
            html += "</div>"
            
        html += "</body></html>"
        return Response(html, mimetype='text/html')
    except Exception as e: return f"Compilation Error: {{e}}", 500


@app.route('/api/graph')
def api_graph():
    with graph_state_lock:
        return jsonify({"nodes": graph_snapshot["nodes"], "edges": graph_snapshot["edges"]})

@app.route('/api/dossier')
def api_dossier():
    node_id = request.args.get('node_id', '').strip()
    url = request.args.get('url', '').strip()

    if node_id:
        with graph_state_lock:
            node = graph_index.get(node_id)
        if node:
            return jsonify({"status": "success", "dossier": node})

    if url and global_dbh:
        record = global_dbh.get_record(url)
        if record:
            return jsonify({"status": "success", "dossier": record})

    return jsonify({"status": "not_found", "dossier": None}), 404

@app.route('/api/task-log')
def api_task_log():
    if not autonomy_executor:
        return jsonify({"tasks": []})
    return jsonify({"tasks": autonomy_executor.get_task_log()})

@app.route('/api/profiles')
def api_profiles():
    if not intel_aggregator:
        return jsonify({"profiles": []})
    return jsonify({"profiles": intel_aggregator.get_all_profiles()})

@app.route('/api/profile/<entity_id>')
def api_profile_detail(entity_id):
    if not intel_aggregator:
        return jsonify({"profile": None}), 404
    profile = intel_aggregator.get_profile(entity_id)
    if not profile:
        return jsonify({"profile": None}), 404
    return jsonify({"profile": profile})

@app.route('/api/tracer/stream')
def tracer_stream():
    # Expects ?seeds=[{"type":"address","value":"...","chain":"ETH"}]
    seeds_param = request.args.get('seeds', '')
    
    if not seeds_param:
        return jsonify({"error": "Missing seeds parameter"}), 400
        
    try:
        seeds = json.loads(seeds_param)
        if not isinstance(seeds, list): raise ValueError
    except:
        return jsonify({"error": "Invalid seeds format. Must be JSON array."}), 400
        
    def generate():
        import asyncio
        yield f"data: {json.dumps({'stage': 'INIT', 'message': f'Initializing NEMESIS Blockchain Tracer for {len(seeds)} inputs'})}\n\n"
        
        if not tracer_engine:
            yield f"data: {json.dumps({'stage': 'ERROR', 'message': 'Tracer Engine failed to initialize or missing dependencies.'})}\n\n"
            return
            
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            
            yield f"data: {json.dumps({'stage': 'TRACING', 'message': 'Reconstructing unified multi-chain graph...'})}\n\n"
            result = loop.run_until_complete(tracer_engine.trace_multiple(seeds))
            
            yield f"data: {json.dumps({'stage': 'INFERENCE', 'message': 'Running Neural Intent Inference on extracted paths...'})}\n\n"
            
            # Save graph snapshot to MongoDB darknet_data
            if global_dbh and global_dbh.dbtype == "mongodb" and global_dbh.coll is not None:
                # Store a unified session graph instead of per-address
                global_dbh.db["tracer_graphs"].insert_one(result)
                yield f"data: {json.dumps({'stage': 'PERSISTENCE', 'message': 'Graph persistently stored to MongoDB (darknet_data/tracer_graphs)'})}\n\n"
            
            yield f"data: {json.dumps({'stage': 'COMPLETE', 'result': result})}\n\n"
            
        except Exception as e:
            yield f"data: {json.dumps({'stage': 'ERROR', 'message': str(e)})}\n\n"
            
    return Response(generate(), mimetype='text/event-stream')

@app.route('/api/documents')
def get_documents():
    if not doc_storage:
        return jsonify({"documents": []})
    return jsonify({"documents": doc_storage.get_storage_log()})

@app.route('/api/stream')
def api_stream():
    def event_stream():
        q = Queue()
        with locks["sse"]: sse_clients.append(q)
        try:
            while True: yield q.get()
        except GeneratorExit:
            with locks["sse"]: sse_clients.remove(q)
    return Response(event_stream(), mimetype="text/event-stream")

@app.route('/api/trigger_report', methods=['POST', 'GET'])
def api_trigger_report():
    success, msg = send_daily_report()
    if success:
        return jsonify({"status": "success", "message": msg})
    else:
        return jsonify({"status": "error", "message": msg}), 500

@app.route('/api/clear_sheet', methods=['GET', 'POST'])
def api_clear_sheet():
    if sheets_manager and sheets_manager.clear_sheet():
        return jsonify({"status": "success", "message": "Google Sheet cleared successfully via Matrix."})
    return jsonify({"status": "error", "message": "Google Sheets API not connected."}), 500

@app.route('/api/flush_sheets', methods=['GET', 'POST'])
def api_flush_sheets():
    if sheets_manager:
        count = sheets_manager.force_flush()
        return jsonify({"status": "success", "message": f"Forced flush of {count} queued entities to Google Sheets."})
    return jsonify({"status": "error", "message": "Google Sheets API not connected."}), 500

@app.route('/api/crawler/start', methods=['POST'])
def api_crawler_start():
    global crawler_is_running, global_worker_stop, global_workers
    if crawler_is_running:
        return jsonify({"status": "error", "message": "Crawler is already running"}), 400
        
    global_worker_stop.clear()
    global_workers = []
    
    # Start the worker threads
    for i in range(MAX_WORKERS):
        t = threading.Thread(target=worker_thread, args=(global_dbh, i+1, global_worker_stop), daemon=True)
        t.start()
        global_workers.append(t)
        
    crawler_is_running = True
    log("[SYSTEM] Crawler threads started via Dashboard.", style="bold green")
    return jsonify({"status": "success", "message": "Crawler started"})

@app.route('/api/crawler/stop', methods=['POST'])
def api_crawler_stop():
    global crawler_is_running, global_worker_stop, global_workers
    if not crawler_is_running:
        return jsonify({"status": "error", "message": "Crawler is not running"}), 400
        
    global_worker_stop.set()
    crawler_is_running = False
    log("[SYSTEM] Crawler threads stopping via Dashboard...", style="bold yellow")
    return jsonify({"status": "success", "message": "Crawler stopping"})

@app.route('/api/crawler/keywords', methods=['GET', 'POST'])
def api_crawler_keywords():
    if request.method == 'GET':
        return jsonify({"status": "success", "keywords": keywords})
        
    data = request.json
    if not data or 'keywords' not in data:
        return jsonify({"status": "error", "message": "No keywords provided"}), 400
        
    new_kws = data['keywords']
    if not isinstance(new_kws, list):
        new_kws = [new_kws]
        
    added = 0
    for kw in new_kws:
        k = kw.strip()
        if k and k not in keywords:
            keywords.append(k)
            added += 1
            
    if added > 0:
        try:
            with open(KEYWORDS_FILE, "a", encoding="utf-8") as f:
                for k in new_kws:
                    f.write(f"{k.strip()}\\n")
        except Exception as e:
            import logging, traceback
            logging.error(f'[Recovered Exception in darknetv2.py] {e}')
            traceback.print_exc()
        log(f"[AUTO] Added {added} new keywords via Dashboard.", style="green")
        
    return jsonify({"status": "success", "message": f"Added {added} keywords", "total": len(keywords)})

@app.route('/api/crawler/status', methods=['GET'])
def api_crawler_status():
    return jsonify({
        "status": "success",
        "running": crawler_is_running,
        "active_threads": len([t for t in global_workers if t.is_alive()]),
        "total_keywords": len(keywords)
    })

def start_api_server(dbh):
    global global_dbh, autonomy_executor, doc_storage, intel_aggregator, sheets_manager
    global_dbh = dbh
    autonomy_executor = AutonomousTaskExecutor(dbh)
    doc_storage = DocumentStorageManager()
    intel_aggregator = IntelligenceProfileAggregator(dbh)
    # Disable reloader to prevent duplicate threads
    app.run(host='0.0.0.0', port=8001, debug=False, use_reloader=False)

# -------------------- DB ABSTRACTION --------------------
class DBHandler:
    def __init__(self, dbtype, conn_info):
        self.dbtype = dbtype.lower()
        self.conn_info = conn_info
        self.client = None
        self.conn = None
        self.cursor = None
        self.db = None
        self.coll = None
        self.lock = threading.Lock()

    def connect_and_prepare(self):
        with self.lock:
            if self.dbtype == "mongodb":
                if MongoClient is None: raise RuntimeError("pymongo not installed")
                try:
                    self.client = MongoClient(self.conn_info, serverSelectionTimeoutMS=8000)
                    self.client.server_info()
                    try: self.db = self.client.get_database()
                    except Exception: self.db = self.client["darknet"]
                    self.coll = self.db["darknet_data"]
                    log(f"MongoDB connected: {self.db.name}.darknet_data (Darknet Collection)", style="green")
                except Exception as e:
                    log(f"MongoDB connection failed: {e}", style="red")
                    self.coll = None
                    
            elif self.dbtype == "sqlite":
                path = self.conn_info.get("path", "crawler_data.db") if isinstance(self.conn_info, dict) else self.conn_info
                self.conn = sqlite3.connect(path, check_same_thread=False)
                self.cursor = self.conn.cursor()
                self._create_sqlite_tables()
                log(f"SQLite connected: {path}", style="green")
            elif self.dbtype == "postgresql":
                if psycopg2 is None: raise RuntimeError("psycopg2 not installed")
                if isinstance(self.conn_info, dict): self.conn = psycopg2.connect(**self.conn_info)
                else: self.conn = psycopg2.connect(self.conn_info)
                self.cursor = self.conn.cursor()
                self._create_postgres_tables()
                log(f"PostgreSQL connected", style="green")
            else:
                log(f"Export-only mode (No DB specified)", style="yellow")

    def _create_sqlite_tables(self):
        create_crawled = """CREATE TABLE IF NOT EXISTS crawled (
            id INTEGER PRIMARY KEY AUTOINCREMENT, hash_id TEXT, crawled_at TEXT, url TEXT UNIQUE, title TEXT,
            description TEXT, query_parameters TEXT, content TEXT, links TEXT, financial TEXT, person TEXT, documents TEXT, keywords_detected TEXT
        );"""
        create_gcei_docs = """CREATE TABLE IF NOT EXISTS gcei_documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT, hash_id TEXT UNIQUE, url TEXT, title TEXT, content TEXT, crawled_at TEXT
        );"""
        create_gcei_entities = """CREATE TABLE IF NOT EXISTS gcei_entities (
            id INTEGER PRIMARY KEY AUTOINCREMENT, entity_id TEXT UNIQUE, entity_type TEXT, value TEXT, confidence REAL, first_seen TEXT
        );"""
        create_gcei_rels = """CREATE TABLE IF NOT EXISTS gcei_relationships (
            id INTEGER PRIMARY KEY AUTOINCREMENT, source_id TEXT, target_id TEXT, rel_type TEXT,
            UNIQUE(source_id, target_id, rel_type)
        );"""
        self.cursor.execute(create_crawled)
        self.cursor.execute(create_gcei_docs)
        self.cursor.execute(create_gcei_entities)
        self.cursor.execute(create_gcei_rels)
        self.conn.commit()

    def _create_postgres_tables(self):
        create_crawled = """CREATE TABLE IF NOT EXISTS crawled (
            id SERIAL PRIMARY KEY, hash_id TEXT, crawled_at TEXT, url TEXT UNIQUE, title TEXT,
            description TEXT, query_parameters TEXT, content TEXT, links TEXT, financial TEXT, person TEXT, documents TEXT, keywords_detected TEXT
        );"""
        create_gcei_docs = """CREATE TABLE IF NOT EXISTS gcei_documents (
            id SERIAL PRIMARY KEY, hash_id TEXT UNIQUE, url TEXT, title TEXT, content TEXT, crawled_at TEXT
        );"""
        create_gcei_entities = """CREATE TABLE IF NOT EXISTS gcei_entities (
            id SERIAL PRIMARY KEY, entity_id TEXT UNIQUE, entity_type TEXT, value TEXT, confidence REAL, first_seen TEXT
        );"""
        create_gcei_rels = """CREATE TABLE IF NOT EXISTS gcei_relationships (
            id SERIAL PRIMARY KEY, source_id TEXT, target_id TEXT, rel_type TEXT,
            UNIQUE(source_id, target_id, rel_type)
        );"""
        self.cursor.execute(create_crawled)
        self.cursor.execute(create_gcei_docs)
        self.cursor.execute(create_gcei_entities)
        self.cursor.execute(create_gcei_rels)
        self.conn.commit()

    def insert_record(self, record):
        try:
            uie_entities = record.get("uie_entities", [])
            fin_data = [e for e in uie_entities if e['ontology_class'] in ['CRYPTO_WALLET', 'BANK_ACCOUNT']]
            per_data = [e for e in uie_entities if e['ontology_class'] in ['PERSON', 'EMAIL', 'PHONE', 'NETWORK_ENTITY']]
            doc_data = [e for e in uie_entities if e['ontology_class'] in ['FILE_DOC', 'ORGANIZATION', 'DOMAIN', 'IP_ADDRESS', 'MALWARE_IOC']]

            # Inject root-level URL to satisfy MongoDB unique indexes
            record["url"] = record["web_info"].get("url")

            if self.dbtype == "mongodb":
                if self.coll is not None:
                    # Upsert prevents Duplicate Key Exceptions on revisits
                    self.coll.update_one({"url": record["url"]}, {"$set": record}, upsert=True)
                else:
                    log(f"[DB] MongoDB collection not initialized, skipping insert for {record['url']}", style="red")
            elif self.dbtype == "sqlite":
                with self.lock:
                    self.cursor.execute(
                        "INSERT OR IGNORE INTO crawled (hash_id,crawled_at,url,title,description,query_parameters,content,links,financial,person,documents,keywords_detected) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)",
                        (
                            record.get("hash-ID"), record.get("crawled_at"), record["web_info"].get("url"),
                            record["web_info"].get("title"), record["web_info"].get("description"),
                            json.dumps(record["web_info"].get("query_parameters", {}), ensure_ascii=False),
                            record["web_info"].get("content"), json.dumps(record["web_info"].get("links", []), ensure_ascii=False),
                            json.dumps(fin_data, ensure_ascii=False),
                            json.dumps(per_data, ensure_ascii=False),
                            json.dumps(doc_data, ensure_ascii=False),
                            json.dumps(record.get("keywords_detected", []), ensure_ascii=False),
                        ),
                    )
                    # Relational Injection
                    self.cursor.execute(
                        "INSERT OR IGNORE INTO gcei_documents (hash_id, url, title, content, crawled_at) VALUES (?,?,?,?,?)",
                        (record.get("hash-ID"), record["web_info"].get("url"), record["web_info"].get("title"), record["web_info"].get("content"), record.get("crawled_at"))
                    )
                    for e in uie_entities:
                        eid = hashlib.sha256(f"{e.get('value')}:{e.get('ontology_class')}".encode()).hexdigest()
                        self.cursor.execute(
                            "INSERT OR REPLACE INTO gcei_entities (entity_id, entity_type, value, confidence, first_seen) VALUES (?,?,?,?,?)",
                            (eid, e.get("ontology_class"), e.get("value"), e.get("confidence", 0.0), e.get("timestamp"))
                        )
                        self.cursor.execute(
                            "INSERT OR IGNORE INTO gcei_relationships (source_id, target_id, rel_type) VALUES (?,?,?)",
                            (record.get("hash-ID"), eid, "EXTRACTED_FROM")
                        )
                    self.conn.commit()
            elif self.dbtype == "postgresql":
                with self.lock:
                    self.cursor.execute(
                        "INSERT INTO crawled (hash_id,crawled_at,url,title,description,query_parameters,content,links,financial,person,documents,keywords_detected) VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s) ON CONFLICT (url) DO NOTHING",
                        (
                            record.get("hash-ID"), record.get("crawled_at"), record["web_info"].get("url"),
                            record["web_info"].get("title"), record["web_info"].get("description"),
                            json.dumps(record["web_info"].get("query_parameters", {}), ensure_ascii=False),
                            record["web_info"].get("content"), json.dumps(record["web_info"].get("links", []), ensure_ascii=False),
                            json.dumps(fin_data, ensure_ascii=False),
                            json.dumps(per_data, ensure_ascii=False),
                            json.dumps(doc_data, ensure_ascii=False),
                            json.dumps(record.get("keywords_detected", []), ensure_ascii=False),
                        ),
                    )
                    self.cursor.execute(
                        "INSERT INTO gcei_documents (hash_id, url, title, content, crawled_at) VALUES (%s,%s,%s,%s,%s) ON CONFLICT (hash_id) DO NOTHING",
                        (record.get("hash-ID"), record["web_info"].get("url"), record["web_info"].get("title"), record["web_info"].get("content"), record.get("crawled_at"))
                    )
                    for e in uie_entities:
                        eid = hashlib.sha256(f"{e.get('value')}:{e.get('ontology_class')}".encode()).hexdigest()
                        self.cursor.execute(
                            "INSERT INTO gcei_entities (entity_id, entity_type, value, confidence, first_seen) VALUES (%s,%s,%s,%s,%s) ON CONFLICT (entity_id) DO UPDATE SET confidence = GREATEST(gcei_entities.confidence, EXCLUDED.confidence)",
                            (eid, e.get("ontology_class"), e.get("value"), e.get("confidence", 0.0), e.get("timestamp"))
                        )
                        self.cursor.execute(
                            "INSERT INTO gcei_relationships (source_id, target_id, rel_type) VALUES (%s,%s,%s) ON CONFLICT DO NOTHING",
                            (record.get("hash-ID"), eid, "EXTRACTED_FROM")
                        )
                    self.conn.commit()
        except Exception as e:
            log(f"[DB] insert error: {e}", style="red")

    def get_record(self, lookup_key):
        try:
            if self.dbtype == "mongodb" and self.coll is not None:
                query = {"$or": [{"url": lookup_key}, {"hash-ID": lookup_key}]}
                doc = self.coll.find_one(query)
                if doc:
                    doc["_id"] = str(doc.get("_id"))
                return doc

            if self.dbtype == "sqlite":
                with self.lock:
                    self.cursor.execute(
                        "SELECT hash_id,crawled_at,url,title,description,query_parameters,content,links,financial,person,documents,keywords_detected FROM crawled WHERE url = ? OR hash_id = ? LIMIT 1",
                        (lookup_key, lookup_key)
                    )
                    row = self.cursor.fetchone()
                if row:
                    return {
                        "hash-ID": row[0],
                        "crawled_at": row[1],
                        "web_info": {
                            "url": row[2],
                            "title": row[3],
                            "description": row[4],
                            "query_parameters": json.loads(row[5] or "{}"),
                            "content": row[6],
                            "links": json.loads(row[7] or "[]")
                        },
                        "uie_entities": json.loads(row[8] or "[]") + json.loads(row[9] or "[]") + json.loads(row[10] or "[]"),
                        "keywords_detected": json.loads(row[11] or "[]")
                    }
            if self.dbtype == "postgresql":
                with self.lock:
                    self.cursor.execute(
                        "SELECT hash_id,crawled_at,url,title,description,query_parameters,content,links,financial,person,documents,keywords_detected FROM crawled WHERE url = %s OR hash_id = %s LIMIT 1",
                        (lookup_key, lookup_key)
                    )
                    row = self.cursor.fetchone()
                if row:
                    return {
                        "hash-ID": row[0],
                        "crawled_at": row[1],
                        "web_info": {
                            "url": row[2],
                            "title": row[3],
                            "description": row[4],
                            "query_parameters": json.loads(row[5] or "{}"),
                            "content": row[6],
                            "links": json.loads(row[7] or "[]")
                        },
                        "uie_entities": json.loads(row[8] or "[]") + json.loads(row[9] or "[]") + json.loads(row[10] or "[]"),
                        "keywords_detected": json.loads(row[11] or "[]")
                    }
        except Exception as e:
            log(f"[DB] query error: {e}", style="red")
        return None

    def close(self):
        try:
            if self.dbtype == "mongodb" and self.client: self.client.close()
            elif self.dbtype in ("sqlite", "postgresql") and self.conn: self.conn.close()
        except Exception: pass

# -------------------- EXTRACTION & PARSING --------------------
class PlaywrightResponse:
    def __init__(self, text, url):
        self.text = text
        self.content = text.encode('utf-8')
        self.status_code = 200
        self.headers = {'Content-Type': 'text/html; charset=utf-8'}
        self.url = url

def safe_get(url, retries=5, timeout=15, allow_redirects=True, stream=False):
    """Fetch URL with exponential backoff, 429 rate-limit awareness, and Playwright JS Fallback."""
    backoff_base = 1.5
    for attempt in range(retries):
        try:
            headers = dict(HEADERS)
            if ua and os.getenv("VITE_CRAWLER_USER_AGENT_ROTATION") == "true":
                headers['User-Agent'] = ua.random
            headers['Accept-Language'] = 'en-US,en;q=0.9'
            headers['Cache-Control'] = 'no-cache'
                
            r = session.get(url, headers=headers, timeout=timeout, allow_redirects=allow_redirects, stream=stream)
            
            # --- HYBRID HEADLESS FALLBACK (PLAYWRIGHT) ---
            # If forbidden or the HTML content is suspiciously small/lacking body content (typical of React/Vue SPA)
            needs_playwright = False
            if not stream:
                if r.status_code in [403, 401]:
                    needs_playwright = True
                elif r.status_code == 200 and len(r.text) < 1500 and ('<noscript>' in r.text or 'enable JavaScript' in r.text):
                    needs_playwright = True

            if needs_playwright and not stream:
                log(f"[HYBRID-JS] Intercepted Anti-Bot / SPA at {url[:50]}. Spawning Headless Chromium...", style="bold magenta")
                try:
                    with sync_playwright() as p:
                        browser = p.chromium.launch(headless=True)
                        context = browser.new_context(user_agent=headers['User-Agent'])
                        page = context.new_page()
                        page.goto(url, timeout=timeout*1000, wait_until="networkidle")
                        html_content = page.content()
                        browser.close()
                        return PlaywrightResponse(html_content, url)
                except Exception as e:
                    log(f"[HYBRID-JS] Headless execution failed: {e}", style="dim red")
                    pass # Fall back to returning original requests response

            if r.status_code == 200:
                return r
            elif r.status_code == 429:  # Rate limited
                retry_after = int(r.headers.get('Retry-After', min(60, backoff_base ** attempt)))
                log(f"[HTTP 429] Rate limited. Backing off {retry_after}s: {url[:50]}...", style="bold yellow")
                time.sleep(retry_after + random.uniform(0, 2))
                continue  # Retry with backoff
            elif r.status_code == 404:
                log(f"[HTTP 404] Not Found: {url}", style="dim")
                return None
            elif r.status_code in [403, 401]:
                log(f"[HTTP {r.status_code}] Forbidden/Unauthorized: {url}", style="dim yellow")
                return None
            else:
                log(f"[HTTP {r.status_code}] Error: {url}", style="dim red")
                if attempt < retries - 1:
                    wait = min(60, backoff_base ** attempt)
                    time.sleep(wait)
                    continue
                return None
        except requests.exceptions.Timeout:
            if attempt < retries - 1:
                wait = min(60, backoff_base ** attempt)
                log(f"[TIMEOUT] Retrying in {wait}s: {url[:50]}...", style="dim yellow")
                time.sleep(wait)
        except requests.exceptions.ConnectionError:
            if attempt < retries - 1:
                wait = min(60, backoff_base ** attempt)
                time.sleep(wait)
        except Exception as e:
            log(f"[EXCEPTION] {str(e)[:80]}: {url[:50]}...", style="dim red")
            time.sleep(min(60, backoff_base ** attempt))
    return None

def handle_query_param_expansion(url):
    parsed = urlparse(url)
    q = parse_qs(parsed.query)
    generated = []
    if not q: return generated
    
    for param in q.keys():
        for kw in list(keywords):
            new_q = {k: v[:] for k, v in q.items()}
            new_q[param] = [kw]
            new_url = parsed._replace(query=urlencode(new_q, doseq=True)).geturl()
            if enqueue_url(new_url):
                generated.append(new_url)
    return generated

def update_stats_from_entity(e_type, ont_class):
    """Helper to update crawler telemetry based on matched entities."""
    if e_type in ['btc', 'bitcoin', 'litecoin', 'dogecoin', 'bitcoin_cash', 'bsv', 'bitcoingold', 'dash', 'zcash']:
        with locks["stats"]: stats["btc_found"] += 1
    elif e_type in ['sol', 'solana']:
        with locks["stats"]: stats["sol_found"] += 1
    elif ont_class == 'CRYPTO_WALLET':
        with locks["stats"]: stats["eth_found"] += 1
    elif e_type == 'email':
        with locks["stats"]: stats["emails_found"] += 1
    elif e_type == 'org':
        with locks["stats"]: stats["orgs_found"] += 1
    elif e_type == 'hash':
        with locks["stats"]: stats["hashes_found"] += 1
    elif e_type == 'documents':
        with locks["stats"]: stats["docs_found"] += 1

def parse_page(url, html):
    soup = BeautifulSoup(html, "lxml")
    text = soup.get_text(" ", strip=True)
    
    metadata = {}
    for meta in soup.find_all("meta"):
        name = meta.get("name", meta.get("property", ""))
        content = meta.get("content", "")
        if name and content and len(content) < 500:
            metadata[name] = content
    
    uie_entities = UIEEngine.extract(text, url)
    uie_entities = enrichment_engine.enrich(uie_entities)
    
    SigmaGraphEngine.stream_nodes(uie_entities) 

    links = []
    for a in soup.find_all("a", href=True):
        href = a.get("href")
        if not href: continue
        href = href.strip()
        try:
            abs_url = urljoin(url, href)
            parsed = urlparse(abs_url)
            if parsed.scheme in ("http", "https") or ".onion" in parsed.netloc:
                links.append(abs_url)
        except Exception: continue
            
    generated = handle_query_param_expansion(url)
    
    # --- HTML Form Search Injection Engine ---
    for form in soup.find_all("form"):
        # Look for search inputs
        search_inputs = []
        for inp in form.find_all("input"):
            name = inp.get("name", "").lower()
            itype = inp.get("type", "").lower()
            if itype == "search" or name in ["q", "query", "search", "keyword", "k"]:
                search_inputs.append(inp)
                
        if search_inputs:
            action = form.get("action", "")
            method = form.get("method", "get").lower()
            form_url = urljoin(url, action)
            
            # Construct submission for each keyword
            for kw in keywords:
                # Gather all other hidden/default inputs
                form_data = {}
                for inp in form.find_all("input"):
                    iname = inp.get("name")
                    if not iname: continue
                    if inp in search_inputs:
                        form_data[iname] = kw
                    else:
                        form_data[iname] = inp.get("value", "")
                
                if method == "get":
                    # For GET, we append to URL and enqueue
                    submit_url = f"{form_url}?{urlencode(form_data)}"
                    if enqueue_url(submit_url):
                        generated.append(submit_url)
                # Note: For POST, we would need to store the POST data in the queue. 
                # Currently enqueue_url only takes a URL. Implementing POST queueing requires queue schema change.
                # We will support GET form injections for now.

    for l in links: enqueue_url(l)
    
    text_lower = text.lower()
    detected_kws = [k for k in keywords if k.lower() in text_lower]
        
    sentiment_result = analyze_sentiment(text[:3000])
    
    record = {
        "hash-ID": sha256(url),
        "crawled_at": datetime.now(timezone.utc).isoformat(),
        "web_info": {
            "url": url,
            "title": soup.title.string if soup.title else "N/A",
            "description": (soup.find("meta", {"name": "description"}) or {}).get("content", "N/A"),
            "metadata": metadata, 
            "query_parameters": parse_qs(urlparse(url).query),
            "content": text[:5000],
            "links": links,
            "subpages": links
        },
        "uie_entities": uie_entities,
        "sentiment_analysis": sentiment_result,
        "keywords_detected": detected_kws
    }
    return record, generated

def send_daily_report():
    smtp_server = os.getenv("SMTP_SERVER", "smtp.gmail.com")
    smtp_port = int(os.getenv("SMTP_PORT", "587"))
    smtp_user = os.getenv("SMTP_USER", "")
    smtp_pass = os.getenv("SMTP_PASS", "")
    to_email = os.getenv("REPORT_EMAIL", "admin@lionsgatenetwork.com")

    if not smtp_user or not smtp_pass:
        log("[REPORT] SMTP credentials not fully configured. Daily report skipped.", style="yellow")
        return False, "SMTP credentials missing"

    try:
        subject = f"Lionsgate NEMESIS Daily OSINT Telemetry Report - {datetime.now().strftime('%Y-%m-%d')}"
        
        with locks["stats"]:
            stats_copy = dict(stats)
        with locks["log"]:
            entities_copy = list(recent_entities)
            
        html_body = f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; color: #333; }}
                h1 {{ color: #0c4a6e; }}
                table {{ border-collapse: collapse; width: 100%; margin-top: 15px; }}
                th, td {{ border: 1px solid #cbd5e1; padding: 8px; text-align: left; }}
                th {{ background-color: #f1f5f9; }}
                .badge {{ padding: 2px 6px; border-radius: 4px; font-weight: bold; font-size: 11px; }}
            </style>
        </head>
        <body>
            <h1>Lionsgate NEMESIS OSINT Report</h1>
            <p>Generated at: {datetime.now().isoformat()}</p>
            
            <h2>System Telemetry</h2>
            <table>
                <tr><th>Metric</th><th>Count</th></tr>
                <tr><td>Processed URLs</td><td>{stats_copy.get('processed', 0)}</td></tr>
                <tr><td>Pending URLs</td><td>{stats_copy.get('pending', 0)}</td></tr>
                <tr><td>Crypto Wallet Hits</td><td>{stats_copy.get('btc_found', 0) + stats_copy.get('eth_found', 0) + stats_copy.get('sol_found', 0)}</td></tr>
                <tr><td>Target Emails Mapped</td><td>{stats_copy.get('emails_found', 0)}</td></tr>
                <tr><td>Forensic Hashes (IOCs)</td><td>{stats_copy.get('hashes_found', 0)}</td></tr>
                <tr><td>Documents Extracted</td><td>{stats_copy.get('docs_found', 0)}</td></tr>
            </table>

            <h2>Recent Mapped Entities</h2>
            <table>
                <thead>
                    <tr>
                        <th>Time</th>
                        <th>Class</th>
                        <th>Value</th>
                        <th>Source URL</th>
                    </tr>
                </thead>
                <tbody>
        """
        for ts, ont, val, url in entities_copy:
            html_body += f"""
                    <tr>
                        <td>{ts}</td>
                        <td><span class="badge">{ont}</span></td>
                        <td><code>{val}</code></td>
                        <td>{url}</td>
                    </tr>
            """
        html_body += """
                </tbody>
            </table>
        </body>
        </html>
        """
        
        msg = MIMEText(html_body, 'html')
        msg['Subject'] = subject
        msg['From'] = smtp_user
        msg['To'] = to_email
        
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(smtp_user, smtp_pass)
            server.sendmail(smtp_user, [to_email], msg.as_string())
            
        log("[REPORT] Daily OSINT report successfully dispatched to " + to_email, style="green")
        return True, "Report successfully sent to " + to_email
    except Exception as e:
        log(f"[REPORT] Failed to send report: {e}", style="red")
        return False, str(e)

# -------------------- AUTOMATED EMAIL REPORTER --------------------
def email_reporter_thread(stop_evt):
    smtp_server = os.getenv("SMTP_SERVER", "smtp.gmail.com")
    smtp_port = int(os.getenv("SMTP_PORT", "587"))
    smtp_user = os.getenv("SMTP_USER", "")
    smtp_pass = os.getenv("SMTP_PASS", "")
    to_email = os.getenv("REPORT_EMAIL", "admin@lionsgatenetwork.com")

    while not stop_evt.is_set():
        now = datetime.now(timezone.utc)
        # Run Cron Job strictly at 00:00 UTC Daily
        if now.hour == 0 and now.minute == 0:
            send_daily_report()
            time.sleep(65) # Sleep past the minute to avoid double triggers
        
        if not smtp_user or not smtp_pass:
            # Supressed repeated missing credential logs in cron mode
            pass

        # Quick check interval allows smooth graceful shutdown checks
        time.sleep(15)

# -------------------- STATE & EXPORTS --------------------
def sha256(text):
    return hashlib.sha256(text.encode("utf-8")).hexdigest()

def init_exports():
    if not os.path.exists(EXPORT_CSV):
        try:
            with open(EXPORT_CSV, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                writer.writerow(["Timestamp", "Ontology_Class", "Entity_Type", "Entity_Value", "Confidence", "Source_URL"])
        except Exception as e:
            log(f"[EXPORT] CSV init error: {e}", style="red")

def save_exports(record):
    try:
        with locks["state"]:
            with open(EXPORT_JSONL, "a", encoding="utf-8") as f:
                f.write(json.dumps(record, ensure_ascii=False, default=str) + "\n")
            
            uie_entities = record.get("uie_entities", [])
            if uie_entities:
                with open(EXPORT_CSV, "a", newline='', encoding="utf-8") as f:
                    writer = csv.writer(f)
                    for e in uie_entities:
                        # Simultaneously export to Google Sheets
                        if sheets_manager:
                            sheets_manager.export_entity(e, record["web_info"].get("url"))
                        writer.writerow([
                            e.get("timestamp"), e.get("ontology_class"), e.get("type"),
                            e.get("value"), e.get("confidence"), e.get("sourceSpan")
                        ])
    except Exception as e:
        log(f"[EXPORT] write error: {e}", style="red")

def save_state_file():
    try:
        with locks["state"]:
            payload = {"queued": list(queued_set), "processed": list(processed_set), "keywords": keywords}
            with open(STATE_FILE, "w", encoding="utf-8") as f: json.dump(payload, f, indent=2)
    except Exception as e: log(f"[STATE] save error: {e}", style="red")

def save_state_periodically(stop_event):
    while not stop_event.is_set():
        time.sleep(SAVE_STATE_INTERVAL)
        save_state_file()

def load_state_file():
    if not os.path.exists(STATE_FILE): return False
    try:
        with open(STATE_FILE, "r", encoding="utf-8") as f: payload = json.load(f)
        q = payload.get("queued", [])
        proc = payload.get("processed", [])
        kw = payload.get("keywords", [])
        for u in q: enqueue_url(u)
        with locks["processed"]: processed_set.update(proc)
        for k in kw:
            if k not in keywords: keywords.append(k)
        log(f"State loaded: Queued={len(q)} Processed={len(proc)}", style="green")
        return True
    except Exception as e:
        log(f"[STATE] load error: {e}", style="red")
        return False

def enqueue_url(u):
    try:
        if not u: return False
        u = u.strip()
        parsed = urlparse(u)
        if not parsed.scheme: return False
        normalized = parsed.geturl().split("#")[0]
        with locks["queued"]:
            if normalized in queued_set or normalized in processed_set: return False
            queue_urls.put(normalized)
            queued_set.add(normalized)
        return True
    except Exception: return False

# -------------------- TOR HANDLING --------------------
def is_tor_running():
    try:
        r = session.get(TOR_CHECK_URL, timeout=8)
        return (r.status_code == 200 and "Congratulations" in r.text)
    except Exception: return False

def start_tor_autorun():
    if not USE_TOR: return True
    if is_tor_running():
        log("[TOR] Already connected.", style="bold green")
        return True
    log("[TOR] Starting Tor process...", style="yellow")
    try:
        if os.name == "nt":
            candidates = ["tor.exe", os.path.join(os.getcwd(), "tor", "tor.exe")]
            for cand in candidates:
                try:
                    subprocess.Popen([cand], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                    time.sleep(1)
                except Exception: continue
        else:
            try: subprocess.Popen(["tor"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            except FileNotFoundError:
                log("[TOR] 'tor' not found. Please install Tor.", style="red")
                return False
    except Exception as e:
        log(f"[TOR] start error: {e}", style="red")
        return False

    for i in range(12):
        time.sleep(2)
        if is_tor_running():
            log("[TOR] Connected securely via port " + TOR_PORT, style="bold green")
            return True
    log("[TOR] Failed to verify connection", style="red")
    return False

# -------------------- GOOGLE SHEETS EXPORT --------------------
def upload_to_google_sheets(record):
    """
    Appends an OSINT record to the target Google Sheet if configured.
    Running in a separate thread so it doesn't block crawler execution.
    """
    def _upload():
        global gs_sheet
        if gs_client is None: return
        
        target_sheet_url = os.getenv("GOOGLE_SHEET_URL")
        if not target_sheet_url: return
        
        try:
            if gs_sheet is None:
                gs_sheet = gs_client.open_by_url(target_sheet_url).sheet1
            
            # Format entities to string
            entities = record.get("uie_entities", [])
            ent_str = "\n".join([f"[{e.get('ontology_class')}] {e.get('value')} ({e.get('confidence')})" for e in entities])
            
            row = [
                record.get("crawled_at", ""),
                record.get("web_info", {}).get("url", ""),
                record.get("web_info", {}).get("title", ""),
                ", ".join(record.get("keywords_detected", [])),
                ent_str
            ]
            gs_sheet.append_row(row)
            log("[GSHEETS] OSINT Data row securely appended to Google Cloud", style="bold green")
        except Exception as e:
            log(f"[GSHEETS] Upload Error: {e}", style="red")
            gs_sheet = None # reset to try connecting again next time
            
    threading.Thread(target=_upload, daemon=True).start()

# -------------------- WORKER THREAD --------------------
def worker_thread(db_handler: DBHandler, tid: int, stop_evt: threading.Event):
    
    # Define document parsing callback inside the worker scope so it can access DB & Locks
    def document_callback(source_url, text_payload, file_meta):
        uie_entities = UIEEngine.extract(text_payload, source_url)
        uie_entities = enrichment_engine.enrich(uie_entities)
        SigmaGraphEngine.stream_nodes(uie_entities)
        
        text_lower = text_payload.lower()
        detected_kws = [k for k in keywords if k.lower() in text_lower]
        
        record = {
            "hash-ID": sha256(source_url),
            "crawled_at": datetime.now(timezone.utc).isoformat(),
            "web_info": {
                "url": source_url,
                "title": file_meta.get("filename", "Binary Document"),
                "description": f"Extracted text from {file_meta.get('type')} file.",
                "metadata": file_meta,
                "query_parameters": parse_qs(urlparse(source_url).query),
                "content": text_payload[:5000],
                "links": [],
                "subpages": []
            },
            "uie_entities": uie_entities,
            "keywords_detected": detected_kws
        }
        
        
        db_handler.insert_record(record)
        save_exports(record)
        upload_to_google_sheets(record)
        
        with locks["processed"]: processed_set.add(source_url)
        with locks["stats"]:
            stats["processed"] += 1
            stats["docs_found"] += 1
            stats["pending"] = queue_urls.qsize()
            
        for e in uie_entities:
            if e['confidence'] >= 0.8:
                ont = e['ontology_class']
                record_entity(ont, e['value'], source_url)
                update_stats_from_entity(e['type'], ont)
                
                # Execute autonomous tasks and build intelligence profile
                if autonomy_executor:
                    entity_id = SigmaGraphEngine.normalize_node_id(e['value'], ont)
                    autonomy_executor.execute_tasks_for_entity(e, source_url)
                    if intel_aggregator:
                        intel_aggregator.build_entity_profile(entity_id, e)
                
        if detected_kws:
            log(f"⚡ Phrases matched in DOC: {detected_kws} at {source_url[:40]}...", style="bold magenta")


    while not stop_evt.is_set():
        try:
            url = queue_urls.get(timeout=2)
        except Empty:
            continue
                
        if url in processed_set:
            queue_urls.task_done()
            continue
            
        # Stream fetch to analyze Content-Type headers before executing heavy loads
        resp = safe_get(url, stream=True)
        if not resp:
            queue_urls.task_done()
            continue
            
        content_type = resp.headers.get('Content-Type', '').lower()
        ext = doc_parser.get_extension_from_url(url)
        document_extensions = [".pdf", ".docx", ".doc", ".pptx", ".ppt", ".csv", ".txt", ".log", ".xlsx", ".xls"]
        
        is_document = False
        if any(x in content_type for x in ['pdf', 'wordprocessing', 'presentation', 'spreadsheet', 'csv', 'excel', 'msword']):
            is_document = True
            if 'pdf' in content_type: ext = '.pdf'
            elif 'csv' in content_type: ext = '.csv'
            elif 'word' in content_type: ext = '.docx'
        elif ext in document_extensions:
            is_document = True
            
        if is_document:
            content_length = int(resp.headers.get('Content-Length', 0))
            if content_length > doc_parser.max_file_size:
                log(f"[-] Skipping {url} - Size exceeds threshold", style="yellow")
            else:
                try:
                    file_bytes = resp.content
                    filename = os.path.basename(urlparse(url).path) or f"doc_{sha256(url)[:8]}{ext}"
                    local_path = os.path.join(doc_parser.download_dir, filename)
                    with open(local_path, "wb") as f:
                        f.write(file_bytes)
                    
                    raw_text = doc_parser.parse_binary_stream(file_bytes, ext)
                    if raw_text and not raw_text.startswith("[PARSING_ERROR]"):
                        document_callback(url, raw_text, {"filename": filename, "type": ext})
                except Exception as e:
                    log(f"[-] Document parsing failed for {url}: {e}", style="red")
        else:
            # Process as standard HTML
            html = resp.text
            
            # Evidentiary HTML Snapshot Preservation
            try:
                snapshot_dir = os.path.join("nemesis_vault", "snapshots")
                os.makedirs(snapshot_dir, exist_ok=True)
                snapshot_path = os.path.join(snapshot_dir, f"snapshot_{sha256(url)}.html")
                with open(snapshot_path, "w", encoding="utf-8") as f:
                    f.write(f"<!-- Captured by Lionsgate NEMESIS OSINT at {datetime.now(timezone.utc).isoformat()} -->\n")
                    f.write(f"<!-- Source: {url} -->\n")
                    f.write(html)
            except Exception as e:
                import logging, traceback
                logging.error(f'[Recovered Exception in darknetv2.py] {e}')
                traceback.print_exc()
                
            try:
                record, generated = parse_page(url, html)
                record["web_info"]["snapshot_path"] = snapshot_path if 'snapshot_path' in locals() else ""
                
                db_handler.insert_record(record)
                save_exports(record)
                upload_to_google_sheets(record)
                
                with locks["processed"]: processed_set.add(url)
                with locks["stats"]:
                    stats["processed"] += 1
                    stats["pending"] = queue_urls.qsize()
                    
                for e in record.get("uie_entities", []):
                    if e['confidence'] >= 0.8:
                        ont = e['ontology_class']
                        record_entity(ont, e['value'], url)
                        update_stats_from_entity(e['type'], ont)
                        
                        if autonomy_executor:
                            entity_id = SigmaGraphEngine.normalize_node_id(e['value'], ont)
                            autonomy_executor.execute_tasks_for_entity(e, url)
                            if intel_aggregator:
                                intel_aggregator.build_entity_profile(entity_id, e)
                    
                if record.get("keywords_detected"):
                    log(f"⚡ Phrases matched: {record['keywords_detected']} at {url[:40]}...", style="bold magenta")
                    
            except Exception as e:
                log(f"[PARSE] error for {url}: {e}", style="red")
                
        queue_urls.task_done()
            
        if stats["processed"] % 15 == 0:
            save_state_file()

# -------------------- RICH LIVE DASHBOARD --------------------
def get_color_for_ontology(ont_class):
    colors = {
        "CRYPTO_WALLET": "yellow",
        "PERSON": "cyan",
        "EMAIL": "magenta",
        "BANK_ACCOUNT": "green",
        "FILE_DOC": "blue",
        "MALWARE_IOC": "red",
        "NETWORK_ENTITY": "bold white",
        "ORGANIZATION": "bold white",
        "DOMAIN": "bold blue"
    }
    return colors.get(ont_class, "white")

def generate_dashboard_layout():
    """Generates the Rich Layout for real-time advanced console rendering."""
    layout = Layout()
    layout.split_column(
        Layout(name="header", size=4),
        Layout(name="main"),
        Layout(name="footer", size=12)
    )
    layout["main"].split_row(
        Layout(name="stats", ratio=1),
        Layout(name="entities", ratio=2)
    )
    
    # Header Panel
    header_text = Align.center(f"[bold cyan]🕸️ LIONSGATE NEMESIS OSINT & DARKNET CRAWLER 🕸️[/bold cyan]\n[bold white]Active Threads: {MAX_WORKERS} | Tor Proxy: {'[green]ENABLED[/green]' if USE_TOR else '[red]DISABLED[/red]'}[/bold white]")
    layout["header"].update(Panel(header_text, box=box.ROUNDED, style="cyan", border_style="cyan"))
    
    # Stats Table
    stats_table = Table(box=box.SIMPLE_HEAVY, expand=True, title="[bold yellow]📊 SYSTEM TELEMETRY[/bold yellow]")
    stats_table.add_column("Metric", style="bold white")
    stats_table.add_column("Count", justify="right", style="bold cyan")
    with locks["stats"]:
        stats_table.add_row("🌐 Pending URLs", f"[yellow]{stats['pending']}[/yellow]")
        stats_table.add_row("✅ Processed URLs", f"[green]{stats['processed']}[/green]")
        stats_table.add_row("💰 Crypto (BTC/ETH)", f"{stats['btc_found']} / {stats['eth_found']}")
        stats_table.add_row("📧 Target Emails", str(stats["emails_found"]))
        stats_table.add_row("🏢 Organizations", str(stats["orgs_found"]))
        stats_table.add_row("🦠 Forensic Hashes", f"[red]{stats['hashes_found']}[/red]")
        stats_table.add_row("📄 Docs Extracted", str(stats["docs_found"]))
    layout["stats"].update(Panel(stats_table, box=box.ROUNDED, border_style="yellow"))
    
    # Live Ontology Table
    entity_table = Table(box=box.MINIMAL_DOUBLE_HEAD, expand=True, title="[bold cyan]📡 REAL-TIME ONTOLOGY EXTRACTION STREAM[/bold cyan]")
    entity_table.add_column("Time", style="dim white", width=10)
    entity_table.add_column("Class", style="bold magenta", width=16)
    entity_table.add_column("Value", style="bold green", width=25)
    entity_table.add_column("Source", style="bold blue")
    
    with locks["log"]:
        # Reverse to show newest at the top
        rows = list(recent_entities)
        for ts, ont, val, url in reversed(rows):
            c = get_color_for_ontology(ont)
            entity_table.add_row(ts, f"[{c}]{ont}[/{c}]", f"[{c}]{val}[/{c}]", url)
            
        # Pad table to maintain stable height and prevent flickering
        for _ in range(15 - len(rows)):
            entity_table.add_row("", "", "", "")

    layout["entities"].update(Panel(entity_table, box=box.ROUNDED, border_style="magenta"))
    
    # Engine Logs
    logs = list(recent_logs)
    while len(logs) < 10:
        logs.append("")
    log_text = "\n".join(logs)
    layout["footer"].update(Panel(log_text, title="[bold green]Engine Subprocess Logs[/bold green]", box=box.ROUNDED, border_style="green"))
    
    return layout

# -------------------- SEEDS & KEYWORDS --------------------
def load_seeds_and_keywords():
    loaded = 0
    if os.path.exists(SEED_FILE):
        with open(SEED_FILE, "r", encoding="utf-8") as f:
            for line in f:
                if enqueue_url(line): loaded += 1
    
    if not AUTONOMOUS_MODE:
        if loaded == 0:
            log("[AUTO] SEED_FILE empty. Enter seed URLs (empty line to finish):", style="bold yellow")
            while True:
                u = ""
                if not u: break
                if enqueue_url(u): loaded += 1
                
        # --- Multi-line, text file, & comma-separated User Keyword Input Prompt ---
        log("[AUTO] Enter custom keywords/phrases (one per line or comma-separated).", style="bold yellow")
        log("[AUTO] You can also type the path to a .txt file to load them.", style="bold yellow")
        log("[AUTO] Press Enter on an empty line to finish.", style="bold yellow")
        
        new_kws = []
        while True:
            user_input = ''
            if not user_input:
                break
            
            if user_input.lower().endswith('.txt') and os.path.isfile(user_input):
                try:
                    with open(user_input, 'r', encoding='utf-8') as f:
                        for line in f:
                            k = line.strip()
                            if k and k not in keywords and k not in new_kws:
                                new_kws.append(k)
                    log(f"[AUTO] Loaded keywords from {user_input}", style="green")
                except Exception as e:
                    log(f"[WARN] Could not read {user_input}: {e}", style="red")
            else:
                for k in user_input.split(","):
                    k = k.strip()
                    if k and k not in keywords and k not in new_kws:
                        new_kws.append(k)
        
        if new_kws:
            keywords.extend(new_kws)
            try:
                with open(KEYWORDS_FILE, "a", encoding="utf-8") as f:
                    for k in new_kws:
                        f.write(f"{k}\n")
            except Exception: pass

    if loaded == 0:
        log("[AUTO] SEED_FILE empty. Hardcoding default entry points...")
        enqueue_url("http://ahmia.fi")
        enqueue_url("https://darknetlive.com")
        enqueue_url("https://news.ycombinator.com")
        enqueue_url("https://pastebin.com")
        enqueue_url("https://reddit.com/r/darknet")
        loaded = 5

    if os.path.exists(KEYWORDS_FILE):
        with open(KEYWORDS_FILE, "r", encoding="utf-8") as f:
            for line in f:
                k = line.strip()
                if k and k not in keywords: keywords.append(k)
                
    if not keywords:
        keywords.extend(DEFAULT_KEYWORDS)
        
    log(f"[AUTO] Tracking {len(keywords)} active keywords.", style="cyan")
    return loaded

# -------------------- AUTO-RUN MAIN --------------------
def main():
    while True:
        try:
            run_crawler_system()
            break  # Exit cleanly if it finishes
        except KeyboardInterrupt:
            log("[SYSTEM] Shutting down gracefully via User Interrupt.", style="bold yellow")
            break
        except Exception as e:
            log(f"[AUTO-HEAL] Critical failure detected: {e}. Rebooting supervisor...", style="bold red")
            time.sleep(5)
            # Supervisor loop continues, auto-healing the process
            pass

def run_crawler_system():
    log("Engine initializing sequence...", style="green")
    
    init_exports()
    if USE_TOR: start_tor_autorun()

    if AUTONOMOUS_MODE:
        log("Autonomous mode active reading from .env", style="green")
        db_type = os.getenv("VITE_CRAWLER_DB_TYPE", "mongodb")
        conn_info = os.getenv("VITE_DATABASE_MONGO_URL", "") if db_type == "mongodb" else os.getenv("VITE_POSTGRES_URI", "crawler_data.db")
    else:
        db_type, conn_info = "sqlite", "crawler_data.db"

    dbh = DBHandler(db_type, conn_info)
    try: dbh.connect_and_prepare()
    except Exception as e: log(f"[DB] init failed: {e}", style="red")

    if not (os.path.exists(STATE_FILE) and load_state_file()):
        load_seeds_and_keywords()

    with locks["stats"]:
        stats["seed"] = queue_urls.qsize()
        stats["pending"] = queue_urls.qsize()
        stats["processed"] = len(processed_set)

    # Launch Flask Backend API in Daemon Thread
    api_thr = threading.Thread(target=start_api_server, args=(dbh,), daemon=True)
    api_thr.start()

    # Spawn secondary UI/Mapper
    spawn_osint_mapper()

    stop_saver = threading.Event()
    saver_thr = threading.Thread(target=save_state_periodically, args=(stop_saver,), daemon=True)
    saver_thr.start()
    
    # Launch Automated Daily Email Reporter
    email_thr = threading.Thread(target=email_reporter_thread, args=(stop_saver,), daemon=True)
    email_thr.start()

    global crawler_is_running, global_worker_stop, global_workers
    global_worker_stop.clear()
    
    if AUTONOMOUS_MODE:
        crawler_is_running = True
        for i in range(MAX_WORKERS):
            t = threading.Thread(target=worker_thread, args=(dbh, i+1, global_worker_stop), daemon=True)
            t.start()
            global_workers.append(t)
    else:
        crawler_is_running = False
        log("[SYSTEM] Crawler waiting for manual start via Dashboard.", style="bold yellow")

    # 🚀 ADVANCED RICH LIVE CONSOLE LOOP
    try:
        with Live(generate_dashboard_layout(), refresh_per_second=4, screen=True) as live:
            while True:
                time.sleep(0.5) # Refresh layout smoothly
                live.update(generate_dashboard_layout())
                if queue_urls.empty():
                    time.sleep(2)
                    if queue_urls.empty(): pass # Do not exit, keep API alive
    except KeyboardInterrupt:
        pass # UI handled shutdown cleanly

    global_worker_stop.set()
    stop_saver.set()
    for t in global_workers: 
        if t.is_alive():
            t.join(timeout=2)
    saver_thr.join(timeout=1)
    save_state_file()
    
    try: dbh.close()
    except Exception: pass

def start_headless_crawler():
    global AUTONOMOUS_MODE
    AUTONOMOUS_MODE = True
    t = threading.Thread(target=main, daemon=True)
    t.start()

if __name__ == "__main__":
    main()